from decimal import Decimal
from html import escape
from io import BytesIO

from django.shortcuts import render, redirect, get_object_or_404

from django.urls import reverse

from django.contrib.auth import authenticate, login, logout

from django.contrib.auth.decorators import login_required

from django.contrib import messages

from django.http import JsonResponse, HttpResponseForbidden, HttpResponse

from django.utils import timezone

from django.db.models import Count, Sum, Q

from django.core.exceptions import PermissionDenied

from datetime import date, timedelta, datetime, time
from django.db.models import Sum

from .models import PayrollAdjustment, RoleMenuPermission
import calendar
import json
from django.contrib.auth.models import Permission
import openpyxl
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


from .models import (

    User, Employee, UserPermission, Department, Designation, Shift, BankAccount,

    Attendance, AttendanceLog, LateEntry, LeaveType, LeaveRequest,

    Notice, Document, DocumentRequest, SalaryStructure, Payslip, PettyCashLedger,

    Asset, Project, Task, OnboardingRecord, SecureFile,

    EmailTemplate, Holiday, NotificationRule, Role,

)

from .forms import (

    LoginForm, EmployeeForm, NoticeForm, LeaveRequestForm,

    AssetForm, SiteSettingsForm, TaskForm, PettyCashForm, DocumentForm, SalaryStructureForm,

    DepartmentForm, DesignationForm, ShiftForm, HolidayForm, NotificationRuleForm,

)

from .decorators import admin_required, manager_or_admin, role_required
from django.views.decorators.http import require_POST
from django.views.decorators.csrf import csrf_exempt
# ─────────────────────────────────────────────
# PERMISSION HELPER
# ─────────────────────────────────────────────

def has_permission(user, module, action="view"):
    """
    Super Admin -> Always True
    Manager / Employee -> Check UserPermission strictly
    """

    # Super Admin always has full access
    if user.role == Role.SUPER_ADMIN:
        return True

    # Check UserPermission for the specific user
    perm = UserPermission.objects.filter(
        user=user,
        module=module
    ).first()

    if not perm:
        return False

    if action == "view":
        return perm.can_view

    if action == "create":
        return perm.can_create

    if action == "edit":
        return perm.can_edit

    if action == "delete":
        return perm.can_delete

    return False

PERMISSION_MODULES = [key for key, _ in UserPermission.MODULE_CHOICES]


def _normalize_role(role):

    if not role:
        return ""

    return role.strip().lower()


# ─── helpers ─────────────────────────────────────────────────────────────────

def _ctx(request, **kw):

    """Common context for every page."""

    view_role = request.GET.get('view_role', request.user.role)

    ctx = {'user': request.user, 'today': date.today(), 'view_role': view_role}

    unread_leaves = LeaveRequest.objects.filter(status='Pending').count()

    ctx['unread_leaves'] = unread_leaves

    # Add allowed menus list to the context using UserPermission
    if request.user.role == 'super_admin':
        ctx['allowed_menus'] = [key for key, _ in UserPermission.MODULE_CHOICES]
    else:
        ctx['allowed_menus'] = list(UserPermission.objects.filter(
            user=request.user,
            can_view=True
        ).values_list('module', flat=True))

    ctx.update(kw)

    return ctx


def _emp(request):

    """Get employee for the logged-in user (or None)."""

    try:
        return request.user.employee
    except Employee.DoesNotExist:
        return None


def _normalize_export_format(export_format):

    value = (export_format or 'xlsx').strip().lower()

    if value == 'ppt':
        value = 'pptx'

    if value not in {'xlsx', 'pdf', 'pptx'}:
        value = 'xlsx'

    return value


def _export_cell_value(value):

    if value is None:
        return ''

    if hasattr(value, 'strftime'):
        try:
            return value.strftime('%Y-%m-%d %H:%M') if hasattr(value, 'hour') else value.strftime('%Y-%m-%d')
        except Exception:
            return str(value)

    return str(value)


def _build_export_sections_payload(sections):

    payload = []

    for section_title, headers, rows in sections:
        payload.append((section_title, list(headers), [list(row) for row in rows]))

    return payload


def _rgb_color(hex_value):

    hex_value = hex_value.lstrip('#')
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def _export_xlsx_response(sections, filename):

    wb = openpyxl.Workbook()

    first_sheet = True

    for section_title, headers, rows in sections:
        ws = wb.active if first_sheet else wb.create_sheet()
        first_sheet = False

        ws.title = (section_title or 'Export')[:31]
        ws.append([_export_cell_value(header) for header in headers])

        if rows:
            for row in rows:
                ws.append([_export_cell_value(value) for value in row])
        else:
            ws.append(['No data available'] + [''] * (max(len(headers) - 1, 0)))

    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    response['Content-Disposition'] = f'attachment; filename="{filename}.xlsx"'
    wb.save(response)
    return response


def _export_pdf_response(title, sections, filename):

    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(letter),
        leftMargin=20,
        rightMargin=20,
        topMargin=24,
        bottomMargin=20,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'ExportTitle',
        parent=styles['Heading1'],
        fontSize=18,
        leading=22,
        textColor=colors.HexColor('#2e3a4f'),
        spaceAfter=8,
    )
    section_style = ParagraphStyle(
        'ExportSection',
        parent=styles['Heading2'],
        fontSize=12,
        leading=14,
        textColor=colors.HexColor('#ef6a4f'),
        spaceAfter=6,
        spaceBefore=8,
    )
    cell_style = ParagraphStyle(
        'ExportCell',
        parent=styles['BodyText'],
        fontSize=8,
        leading=10,
        wordWrap='CJK',
    )

    elements = [Paragraph(escape(title), title_style)]

    for index, (section_title, headers, rows) in enumerate(sections):
        if index > 0:
            elements.append(Spacer(1, 8))

        elements.append(Paragraph(escape(section_title), section_style))

        table_rows = [headers] + rows if rows else [headers, ['No data available'] + [''] * (max(len(headers) - 1, 0))]
        table_data = []

        for row in table_rows:
            table_data.append([
                Paragraph(escape(_export_cell_value(value)).replace('\n', '<br/>'), cell_style)
                for value in row
            ])

        column_count = max(len(headers), 1)
        usable_width = doc.width
        column_width = usable_width / column_count
        table = Table(table_data, colWidths=[column_width] * column_count, repeatRows=1)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#ef6a4f')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 8),
            ('LEADING', (0, 0), (-1, -1), 10),
            ('GRID', (0, 0), (-1, -1), 0.35, colors.HexColor('#d9dee7')),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.whitesmoke, colors.HexColor('#f8fafc')]),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 6),
            ('RIGHTPADDING', (0, 0), (-1, -1), 6),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
        ]))
        elements.append(table)

    doc.build(elements)

    response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{filename}.pdf"'
    return response


def _export_pptx_response(title, sections, filename):

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = 'Export generated from Luminous HRM'

    blank_layout = prs.slide_layouts[6]

    for section_title, headers, rows in sections:
        rows = rows or [['No data available'] + [''] * (max(len(headers) - 1, 0))]
        chunk_size = 8

        for offset in range(0, len(rows), chunk_size):
            chunk = rows[offset:offset + chunk_size]
            slide = prs.slides.add_slide(blank_layout)

            title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.2), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = f'{title} - {section_title}'
            title_frame.paragraphs[0].font.size = Pt(22)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = _rgb_color('#2e3a4f')

            rows_count = len(chunk) + 1
            cols_count = max(len(headers), 1)
            table = slide.shapes.add_table(
                rows_count,
                cols_count,
                Inches(0.4),
                Inches(1.1),
                Inches(12.4),
                Inches(5.7),
            ).table

            for col_index, header in enumerate(headers):
                cell = table.cell(0, col_index)
                cell.text = _export_cell_value(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb_color('#ef6a4f')
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(10)
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.alignment = PP_ALIGN.CENTER

            for row_index, row in enumerate(chunk, start=1):
                for col_index, value in enumerate(row):
                    cell = table.cell(row_index, col_index)
                    cell.text = _export_cell_value(value)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(9)
                        paragraph.font.color.rgb = _rgb_color('#2e3a4f')

            for row in table.rows:
                for cell in row.cells:
                    cell.text_frame.word_wrap = True

    buffer = BytesIO()
    prs.save(buffer)
    response = HttpResponse(
        buffer.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = f'attachment; filename="{filename}.pptx"'
    return response


def _export_tabular_response(title, sections, filename, export_format):

    export_format = _normalize_export_format(export_format)
    sections = _build_export_sections_payload(sections)

    if export_format == 'pdf':
        return _export_pdf_response(title, sections, filename)

    if export_format == 'pptx':
        return _export_pptx_response(title, sections, filename)

    return _export_xlsx_response(sections, filename)


@csrf_exempt
def login_view(request):

    if request.user.is_authenticated:
        return redirect('dashboard')

    form = LoginForm(request, data=request.POST or None)

    if request.method == 'POST' and form.is_valid():
        user = form.get_user()
        login(request, user)
        return redirect('dashboard')

    return render(request, 'hrm/login.html', {'form': form})


def logout_view(request):

    logout(request)
    return redirect('login')


@login_required
def dashboard(request):

    view_role = request.GET.get('view_role', request.user.role)
    emp = _emp(request)
    today = date.today()

    emp_count = Employee.objects.count()
    present = Attendance.objects.filter(date=today, status__in=['Present', 'Late']).count()
    pending_lv = LeaveRequest.objects.filter(from_date__lte=today, to_date__gte=today, status='Approved').count()
    late_today = Attendance.objects.filter(date=today, status='Late').count()
    total_employees = Employee.objects.count()
    # ==================== TREND CALCULATIONS ====================
    from dateutil.relativedelta import relativedelta   # make sure this import is at the top

    this_month_start = today.replace(day=1)
    last_month_start = this_month_start - relativedelta(months=1)
    last_month_end = this_month_start - timedelta(days=1)
    yesterday = today - timedelta(days=1)

    # 1. Employees Trend (new hires this month vs last month)
    this_month_emp = Employee.objects.filter(join_date__gte=this_month_start).count()
    last_month_emp = Employee.objects.filter(
        join_date__gte=last_month_start,
        join_date__lte=last_month_end
    ).count()

    if last_month_emp > 0:
        emp_change = round(((this_month_emp - last_month_emp) / last_month_emp) * 100)
    else:
        emp_change = this_month_emp

    employees_trend = f"+{emp_change}" if emp_change >= 0 else str(emp_change)
    employees_trend_dir = "up" if emp_change >= 0 else "down"

    # 2. Present Today Trend (today vs yesterday)
    present_yesterday = Attendance.objects.filter(
        date=yesterday, status__in=['Present', 'Late']
    ).count()

    if present_yesterday > 0:
        present_change = round(((present - present_yesterday) / present_yesterday) * 100)
    else:
        present_change = 0

    present_trend = f"{present_change}%" if present_change >= 0 else f"{present_change}%"
    present_trend_dir = "up" if present_change >= 0 else "down"

    # 3. On Leave Trend (today vs 7 days ago)
    last_week = today - timedelta(days=7)
    leave_last_week = LeaveRequest.objects.filter(
        from_date__lte=last_week, to_date__gte=last_week, status='Approved'
    ).count()

    leave_change = pending_lv - leave_last_week
    leave_trend = f"+{leave_change}" if leave_change >= 0 else str(leave_change)
    leave_trend_dir = "up" if leave_change >= 0 else "down"

    # 4. Payroll Trend (you can replace with real Salary calculation later)
    # Temporary static for now – replace when you have Salary model
    payroll_trend = "+4%"
    payroll_trend_dir = "up"
    recent_notices = Notice.objects.all()[:5]
    recent_tasks = Task.objects.filter(Q(assignee=emp) | Q(project__manager=emp)).order_by('-created_at')[:5] if emp else Task.objects.all()[:5]

    chart = []
    for i in range(5, -1, -1):
        m = (today.replace(day=1) - timedelta(days=i * 28)).replace(day=1)
        cnt = Attendance.objects.filter(date__year=m.year, date__month=m.month, status='Present').count()
        chart.append({'month': m.strftime('%b'), 'count': cnt})

    team_size = 0
    team_present = 0
    pending_approvals = 0
    team_tasks = 0
    days_present = 0
    late_count = 0
    open_tasks = 0

    if view_role == 'manager' or request.user.role == Role.MANAGER:
        team_size = Employee.objects.filter(department=emp.department if emp else None, status='Active').count() if emp else 0
        team_present = Attendance.objects.filter(date=today, employee__department=emp.department if emp else None, status__in=['Present', 'Late']).count() if emp else 0
        pending_approvals = LeaveRequest.objects.filter(status='Pending').count()
        team_tasks = Task.objects.filter(assignee__department=emp.department if emp else None, status='Pending').count() if emp else 0
    elif view_role == 'employee' or request.user.role == Role.EMPLOYEE:
        days_present = Attendance.objects.filter(employee=emp, date__month=today.month, date__year=today.year, status='Present').count() if emp else 0
        late_count = Attendance.objects.filter(employee=emp, date__month=today.month, date__year=today.year, status='Late').count() if emp else 0
        open_tasks = Task.objects.filter(assignee=emp, status='Pending').count() if emp else 0

    casual_total = 3
    sick_total = 3
    special_total = 2
    casual_balance = casual_total
    sick_balance = sick_total
    special_balance = special_total

    if emp:
        casual_used = LeaveRequest.objects.filter(employee=emp, leave_type__name='Casual', status='Approved').count()
        sick_used = LeaveRequest.objects.filter(employee=emp, leave_type__name='Sick', status='Approved').count()
        special_used = LeaveRequest.objects.filter(employee=emp, leave_type__name='Special', status='Approved').count()
        casual_balance = casual_total - casual_used
        sick_balance = sick_total - sick_used
        special_balance = special_total - special_used

    departments = Department.objects.annotate(employee_count=Count('employee')).filter(employee_count__gt=0)
    workforce_data = []
    for dept in departments:
        percentage = (dept.employee_count / total_employees * 100) if total_employees > 0 else 0
        workforce_data.append({'name': dept.name, 'count': dept.employee_count, 'percentage': round(percentage, 1)})

    return render(request, 'hrm/dashboard.html', _ctx(
        request,
        emp_count=emp_count,
        present=present,
        pending_lv=pending_lv,
        late_today=late_today,
        recent_notices=recent_notices,
        recent_tasks=recent_tasks,
        chart=json.dumps(chart),
        view_role=view_role,
        team_size=team_size,
        team_present=team_present,
        pending_approvals=pending_approvals,
        team_tasks=team_tasks,
        days_present=days_present,
        late_count=late_count,
        open_tasks=open_tasks,
        casual_balance=casual_balance,
        casual_total=casual_total,
        sick_balance=sick_balance,
        sick_total=sick_total,
        special_balance=special_balance,
        special_total=special_total,
        workforce_data=workforce_data,
        total_employees=total_employees,
        # Dynamic trends
        employees_trend=employees_trend,
        employees_trend_dir=employees_trend_dir,
        present_trend=present_trend,
        present_trend_dir=present_trend_dir,
        leave_trend=leave_trend,
        leave_trend_dir=leave_trend_dir,
        payroll_trend=payroll_trend,
        payroll_trend_dir=payroll_trend_dir,        
    ))


@login_required
def notice_list(request):
    if not has_permission(request.user, "notice", "view"):
        messages.error(request, "You don't have permission to access notices.")
        return redirect("dashboard")

    from django.utils import timezone

    notices = Notice.objects.all()

    now = timezone.now()

    pinned_count = notices.filter(pinned=True).count()

    this_month_count = notices.filter(
        created_at__year=now.year,
        created_at__month=now.month
    ).count()

    policy_count = notices.filter(
        category="Policy"
    ).count()

    return render(
        request,
        "hrm/notice_list.html",
        _ctx(
            request,
            notices=notices,
            pinned_count=pinned_count,
            this_month_count=this_month_count,
            policy_count=policy_count,
        )
    )


@login_required
def notice_detail(request, pk):
    if not has_permission(request.user, "notice", "view"):
        messages.error(request, "You don't have permission to view notice.")
        return redirect("notice_list")

    notice = get_object_or_404(
        Notice,
        pk=pk
    )

    return render(
        request,
        "hrm/notice_detail.html",
        _ctx(
            request,
            notice=notice
        )
    )


@manager_or_admin
def notice_create(request):
    if not has_permission(request.user, "notice", "create"):
        messages.error(request, "You don't have permission to create notice.")
        return redirect("notice_list")
    

    if request.method == "POST":

        Notice.objects.create(
            title=request.POST.get("title"),
            body=request.POST.get("body"),
            category=request.POST.get("category"),
            audience=request.POST.get("audience", "all"),
            posted_by=request.user,
            pinned=bool(request.POST.get("pinned")),
        )


        messages.success(
            request,
            "Notice created successfully."
        )


        return redirect("notice_list")


    return render(
        request,
        "hrm/notice_form.html",
        _ctx(request)
    )


@login_required
def notice_export(request):
    if not has_permission(request.user, "notice", "view"):
        messages.error(request, "You don't have permission to export notices.")
        return redirect("notice_list")

    notices = Notice.objects.all()

    sections = [
        (
            'Notices',
            ['Title', 'Category', 'Message', 'Posted By', 'Date'],
            [
                [
                    n.title,
                    n.category,
                    n.body,
                    n.posted_by.get_full_name() if n.posted_by else 'HR Department',
                    n.created_at.strftime('%Y-%m-%d'),
                ]
                for n in notices
            ],
        )
    ]

    return _export_tabular_response('Notice Export', sections, 'Notice', request.GET.get('format'))


@login_required
def notice_notify(request, pk):
    if not has_permission(request.user, "notice", "edit"):
        messages.error(request, "You don't have permission to resend notice.")
        return redirect("notice_list")

    notice = get_object_or_404(
        Notice,
        pk=pk
    )


    messages.success(
        request,
        "Notice resent to all employees."
    )


    return redirect(
        "notice_detail",
        pk=notice.pk
    )


@manager_or_admin
def notice_delete(request, pk):
    if not has_permission(request.user, "notice", "delete"):
        messages.error(request, "You don't have permission to delete notice.")
        return redirect("notice_list")

    notice = get_object_or_404(
        Notice,
        pk=pk
    )

    notice.delete()

    messages.success(
        request,
        "Notice deleted successfully."
    )

    return redirect("notice_list")
# ─── EMPLOYEES ────────────────────────────────────────────────────────────────

@login_required

def user_list(request):
    if not has_permission(request.user, "users", "view"):
        messages.error(request, "You don't have permission to access employees.")
        return redirect("dashboard")

    q    = request.GET.get('q', '')

    dept = request.GET.get('dept', '')

    emps = Employee.objects.select_related('user', 'department', 'designation').all()

    if q:

        emps = emps.filter(

            Q(user__first_name__icontains=q) | Q(user__last_name__icontains=q) |

            Q(user__email__icontains=q) | Q(emp_id__icontains=q)

        )

    if dept:

        emps = emps.filter(department_id=dept)



    # managers/employees only see their dept

    if request.user.role == Role.EMPLOYEE:

        if _emp(request):

            emps = emps.filter(department=_emp(request).department)



    departments = Department.objects.all()

    return render(request, 'hrm/user_list.html', _ctx(

        request, employees=emps, departments=departments, q=q, dept=dept

    ))





@login_required
def user_detail(request, pk):
    emp = get_object_or_404(Employee, pk=pk)

    # Employee can only see own profile
    if request.user.role == Role.EMPLOYEE:
        if emp.user != request.user:
            raise PermissionDenied

    # My Profile highlight
    is_my_profile = (request.user == emp.user)

    leaves = LeaveRequest.objects.filter(employee=emp)[:5]
    assets = Asset.objects.filter(assigned_to=emp)
    tasks = Task.objects.filter(assignee=emp)[:5]
    payslips = Payslip.objects.filter(employee=emp)[:3]

    return render(request, 'hrm/user_detail.html', _ctx(
        request,
        emp=emp,
        leaves=leaves,
        assets=assets,
        tasks=tasks,
        payslips=payslips,
        is_my_profile=is_my_profile
    ))
@admin_required

def user_create(request):
    if not has_permission(request.user, "users", "create"):
        messages.error(request, "You don't have permission to create user.")
        return redirect("user_list")

    form = EmployeeForm(request.POST or None)

    if form.is_valid():

        # Create User

        email = form.cleaned_data['email']

        u = User.objects.create_user(

            username=email,

            email=email,

            first_name=form.cleaned_data['first_name'],

            last_name=form.cleaned_data['last_name'],

            password='12345',

            role=form.cleaned_data['role'],

        )

        emp = form.save(commit=False)

        emp.user = u

        emp.save()

        # Initialize default permissions for new user based on role
        from .models import UserPermission, RoleMenuPermission
        
        if u.role != 'super_admin':
            # Copy role default permissions to user-specific permissions
            role_perms = RoleMenuPermission.objects.filter(role=u.role)
            for role_perm in role_perms:
                UserPermission.objects.create(
                    user=u,
                    module=role_perm.menu,
                    can_view=role_perm.can_view,
                    can_create=role_perm.can_create,
                    can_edit=role_perm.can_edit,
                    can_delete=role_perm.can_delete,
                )

        messages.success(request, f'Employee {u.get_full_name()} created. Default password: 12345')

        return redirect('user_list')

    return render(request, 'hrm/form.html', _ctx(request, form=form, title='Add Employee', back='user_list'))


@login_required
def employee_export(request):
    if not has_permission(request.user, "users", "view"):
        messages.error(request, "You don't have permission to export employees.")
        return redirect("user_list")

    employees = Employee.objects.select_related('user', 'department', 'designation').all()

    sections = [
        (
            'Employees',
            ['Employee ID', 'Name', 'Email', 'Role', 'Department', 'Designation', 'Status', 'Join Date'],
            [
                [
                    emp.emp_id,
                    emp.user.get_full_name() or emp.user.username,
                    emp.user.email,
                    emp.user.get_role_display(),
                    emp.department.name if emp.department else '',
                    emp.designation.title if emp.designation else '',
                    emp.status,
                    emp.join_date.strftime('%Y-%m-%d') if emp.join_date else '',
                ]
                for emp in employees
            ],
        )
    ]

    return _export_tabular_response('Employee Export', sections, 'Employees', request.GET.get('format'))





@admin_required

def user_edit(request, pk):
    if not has_permission(request.user, "users", "edit"):
        messages.error(request, "You don't have permission to edit user.")
        return redirect("user_list")

    emp  = get_object_or_404(Employee, pk=pk)

    form = EmployeeForm(request.POST or None, instance=emp)

    if form.is_valid():

        emp.user.first_name = form.cleaned_data['first_name']

        emp.user.last_name  = form.cleaned_data['last_name']

        emp.user.email      = form.cleaned_data['email']

        emp.user.role       = form.cleaned_data['role']

        emp.user.save()

        form.save()

        messages.success(request, 'Employee updated.')

        return redirect('user_detail', pk=pk)

    return render(request, 'hrm/form.html', _ctx(request, form=form, title='Edit Employee', back='user_list'))





@admin_required

def user_delete(request, pk):
    if not has_permission(request.user, "users", "delete"):
        messages.error(request, "You don't have permission to delete user.")
        return redirect("user_list")

    emp = get_object_or_404(Employee, pk=pk)

    emp.user.delete()

    messages.success(request, 'Employee deleted.')

    return redirect('user_list')

@login_required
def user_permissions(request, pk=None):

    # Only Super Admin can manage permissions
    if request.user.role != Role.SUPER_ADMIN:
        messages.error(
            request,
            "You don't have permission to access this page."
        )
        return redirect("user_list")

    # ─────────────────────────────────────────────
    # MENU LIST
    # ─────────────────────────────────────────────

    MENU_CHOICES = UserPermission.MODULE_CHOICES

    # ─────────────────────────────────────────────
    # SELECT USER
    # ─────────────────────────────────────────────

    selected_user = None

    if pk is not None:
        selected_user = get_object_or_404(User, pk=pk)

    # Role selected from dropdown
    selected_role = (
        request.GET.get("role", "")
        or (selected_user.role if selected_user else "")
    )

    # ─────────────────────────────────────────────
    # SAVE ROLE DEFAULT PERMISSIONS
    # ─────────────────────────────────────────────

    if request.method == "POST":

        permission_type = request.POST.get("permission_type")

        # =====================================================
        # 1. SAVE ROLE DEFAULT PERMISSIONS
        # =====================================================

        if permission_type == "role":

            role = request.POST.get("role")

            if not role:
                messages.error(request, "Please select a role first.")
                return redirect("user_permissions")

            for module, module_name in MENU_CHOICES:

                RoleMenuPermission.objects.update_or_create(
                    role=role,
                    menu=module,
                    defaults={
                        "can_view": request.POST.get(
                            f"{module}_view"
                        ) == "on",

                        "can_create": request.POST.get(
                            f"{module}_create"
                        ) == "on",

                        "can_edit": request.POST.get(
                            f"{module}_edit"
                        ) == "on",

                        "can_delete": request.POST.get(
                            f"{module}_delete"
                        ) == "on",
                    }
                )

            messages.success(
                request,
                f"Default permissions for {role} updated successfully."
            )

            return redirect(
                f"{request.path}?role={role}"
            )

        # =====================================================
        # 2. SAVE INDIVIDUAL USER PERMISSIONS
        # =====================================================

        elif permission_type == "user":

            if selected_user is None:
                messages.error(
                    request,
                    "Please select a user first."
                )
                return redirect("user_permissions")

            # Delete all existing permissions for this user first
            UserPermission.objects.filter(user=selected_user).delete()

            # Only create permissions for modules that have at least one checkbox checked
            for module, module_name in MENU_CHOICES:
                has_any_permission = (
                    request.POST.get(f"{module}_view") == "on" or
                    request.POST.get(f"{module}_create") == "on" or
                    request.POST.get(f"{module}_edit") == "on" or
                    request.POST.get(f"{module}_delete") == "on"
                )
                
                if has_any_permission:
                    UserPermission.objects.create(
                        user=selected_user,
                        module=module,
                        can_view=request.POST.get(f"{module}_view") == "on",
                        can_create=request.POST.get(f"{module}_create") == "on",
                        can_edit=request.POST.get(f"{module}_edit") == "on",
                        can_delete=request.POST.get(f"{module}_delete") == "on",
                    )

            messages.success(
                request,
                f"Permissions for {selected_user.get_full_name() or selected_user.username} updated successfully."
            )

            return redirect(
                f"/employees/{selected_user.pk}/permissions/"
            )

    # ─────────────────────────────────────────────
    # ROLE DEFAULT PERMISSIONS
    # ─────────────────────────────────────────────

    role_permissions = {}

    if selected_role:

        role_permissions_qs = RoleMenuPermission.objects.filter(
            role=selected_role
        )

        role_permissions = {
            item.menu: item
            for item in role_permissions_qs
        }

    # ─────────────────────────────────────────────
    # INDIVIDUAL USER PERMISSIONS
    # ─────────────────────────────────────────────

    user_permissions_data = {}

    if selected_user:

        user_permissions_qs = UserPermission.objects.filter(
            user=selected_user
        )

        user_permissions_data = {
            item.module: item
            for item in user_permissions_qs
        }

    # ─────────────────────────────────────────────
    # BUILD MENU DATA FOR TEMPLATE
    # ─────────────────────────────────────────────

    permission_rows = []

    for module, module_name in MENU_CHOICES:

        role_perm = role_permissions.get(module)

        user_perm = user_permissions_data.get(module)

        permission_rows.append({
            "module": module,
            "name": module_name,

            "role_can_view": (
                role_perm.can_view
                if role_perm else False
            ),

            "role_can_create": (
                role_perm.can_create
                if role_perm else False
            ),

            "role_can_edit": (
                role_perm.can_edit
                if role_perm else False
            ),

            "role_can_delete": (
                role_perm.can_delete
                if role_perm else False
            ),

            "user_can_view": (
                user_perm.can_view
                if user_perm else None
            ),

            "user_can_create": (
                user_perm.can_create
                if user_perm else None
            ),

            "user_can_edit": (
                user_perm.can_edit
                if user_perm else None
            ),

            "user_can_delete": (
                user_perm.can_delete
                if user_perm else None
            ),
        })

    # ─────────────────────────────────────────────
    # RENDER
    # ─────────────────────────────────────────────

    return render(
        request,
        "hrm/user_permissions.html",
        _ctx(
            request,

            selected_user=selected_user,
            selected_role=selected_role,

            menu_choices=MENU_CHOICES,

            role_permissions=role_permissions,
            user_permissions_data=user_permissions_data,

            permission_rows=permission_rows,
        )
    )
    # ── POST ──────────────────────────────────────────────────────────────────
    if request.method == "POST":
        # A) ONE USER only
        user_id = request.POST.get("user_id")
        if user_id:
            target = get_object_or_404(User, pk=user_id)
            role_for_menus = target.role
            menu_choices = menus_for_role(role_for_menus)

            # Module permissions — ONLY this user
            for menu_key, _ in menu_choices:
                UserPermission.objects.update_or_create(
                    user=target,
                    module=menu_key,
                    defaults={
                        "can_view": request.POST.get(f"{menu_key}_view") == "on",
                        "can_create": request.POST.get(f"{menu_key}_create") == "on",
                        "can_edit": request.POST.get(f"{menu_key}_edit") == "on",
                        "can_delete": request.POST.get(f"{menu_key}_delete") == "on",
                    },
                )

            messages.success(
                request,
                f"Permissions updated for {target.get_full_name()} only. Other users were not changed."
            )
            return redirect("user_permissions", pk=target.pk)

        # B) ROLE defaults only (does NOT loop every user)
        role = request.POST.get("role", "")
        if not role:
            messages.error(request, "Please select a role first.")
            return redirect("user_permissions")

        menu_choices = menus_for_role(role)

        RoleMenuPermission.objects.filter(role=role).delete()
        for menu_key, _ in menu_choices:
            if request.POST.get(f"menu_{menu_key}"):
                RoleMenuPermission.objects.create(
                    role=role,
                    menu=menu_key,
                    is_allowed=True,
                )

        messages.success(
            request,
            f"Default menu permissions for role '{role}' saved. "
            "Individual user custom permissions were not changed."
        )
        return redirect(f"/employees/permissions/?role={role}")

    # ── GET display ───────────────────────────────────────────────────────────
    saved_permissions = {}
    allowed_menus = []

    if selected_user:
        saved_permissions = {
            p.module: p
            for p in UserPermission.objects.filter(user=selected_user)
        }
        # Menus for display: from role defaults (user-specific module flags still apply in has_permission)
        allowed_menus = list(
            RoleMenuPermission.objects.filter(
                role=selected_user.role,
                is_allowed=True,
            ).values_list("menu", flat=True)
        )
        # If user has any can_view, prefer those for menu checkboxes
        user_view_menus = [
            p.module for p in saved_permissions.values() if p.can_view
        ]
        if user_view_menus:
            allowed_menus = user_view_menus
    elif selected_role:
        allowed_menus = list(
            RoleMenuPermission.objects.filter(
                role=selected_role,
                is_allowed=True,
            ).values_list("menu", flat=True)
        )

    role_user_count = (
        User.objects.filter(role=selected_role).count() if selected_role else 0
    )
    selected_role_label = dict(Role.choices).get(selected_role, "") if selected_role else ""

    return render(request, "hrm/user_permissions.html", {
        "saved_permissions": saved_permissions,
        "selected_user": selected_user,
        "selected_role": selected_role,
        "selected_role_label": selected_role_label,
        "role_choices": Role.choices,
        "role_user_count": role_user_count,
        "allowed_menus": allowed_menus,
        "menu_choices": menu_choices,
    })
# ─── MY PROFILE ───────────────────────────────────────────────────────────────
@login_required
def my_profile(request):
    emp = _emp(request)
    if not emp:
        messages.error(request, "No employee profile found.")
        return redirect('dashboard')
    
    # Get the view role from request parameter or use user's actual role
    view_role = request.GET.get('view_role', request.user.role)
    
    # Always go to own detail page with view_role preserved
    return redirect(f'/employees/{emp.pk}/?view_role={view_role}')




# ─── ATTENDANCE ───────────────────────────────────────────────────────────────
@login_required
def attendance_list(request):
    if not has_permission(request.user, "attendance", "view"):
        messages.error(request, "You don't have permission to access Attendance.")
        return redirect("dashboard")
    today = date.today()

    month = int(request.GET.get('month', today.month))
    year = int(request.GET.get('year', today.year))

    view_mode = request.GET.get('view', 'list')

    today_attendance = None


    # Current employee
    current_employee = _emp(request)

    if current_employee:
        today_attendance = Attendance.objects.filter(
            employee=current_employee,
            date=today
        ).first()


    if request.user.role == Role.EMPLOYEE:

        employee = current_employee

        if employee:


            records = Attendance.objects.filter(
                employee=employee,
                date__month=month,
                date__year=year
            ).select_related(
                'employee__user'
            ).order_by('-date')


            employees = Employee.objects.filter(
                pk=employee.pk
            ).select_related(
                'user'
            )


        else:

            records = Attendance.objects.none()
            employees = Employee.objects.none()


    else:

        records = Attendance.objects.filter(
            date__month=month,
            date__year=year
        ).select_related(
            'employee__user'
        ).order_by('-date')


        employees = Employee.objects.select_related(
            'user'
        ).order_by(
            'user__first_name',
            'user__last_name'
        )



    # Today's attendance
    today_records = Attendance.objects.filter(
        date=today
    ).select_related(
        'employee__user'
    )



    attendance_by_employee = {
        record.employee_id: record
        for record in today_records
    }



    approved_leaves = LeaveRequest.objects.filter(
        status='Approved',
        from_date__lte=today,
        to_date__gte=today
    ).values_list(
        'employee_id',
        flat=True
    )



    present_count = sum(
        1 for record in today_records
        if record.status == 'Present'
    )


    late_count = sum(
        1 for record in today_records
        if record.status == 'Late'
    )


    on_leave_count = len(set(approved_leaves))


    absent_count = 0


    attendance_rows = []



    for employee in employees:


        record = attendance_by_employee.get(
            employee.pk
        )


        if record:


            if record.status == "Late":

                status_label = "Late"
                status_class = "amber"


            elif record.status == "Present":

                status_label = "Present"
                status_class = "green"


            elif record.status == "Leave":

                status_label = "On Leave"
                status_class = "gray"


            else:

                status_label = record.status
                status_class = "gray"



        else:

            status_label = "Absent"
            status_class = "red"

            absent_count += 1



        attendance_rows.append({

            "employee": employee,

            "check_in": record.check_in if record else None,

            "check_out": record.check_out if record else None,

            "status_label": status_label,

            "status_class": status_class,

        })



    # Punch logs
    punch_logs = AttendanceLog.objects.filter(
        date=today
    ).select_related(
        'employee__user'
    )[:20]



    # Manual attendance entries
    manual_entries = Attendance.objects.filter(
        date=today
    ).select_related(
        'employee__user'
    )



    # Device / sync data
    device_names = []

    devices_online = 0

    punches_synced = punch_logs.count()

    missing_punches = 0

    last_sync_label = "Just now"



    active_tab = request.GET.get('tab', 'today')

    if active_tab not in {'today', 'manual', 'sync', 'reports'}:

        active_tab = 'today'
    report_type = request.GET.get(
       'report_type',
    'summary'
)
    show_absent_report = report_type == "absent"

    show_overtime_report = report_type == "overtime"

    show_summary_report = report_type == "summary"
    


    report_departments = Department.objects.annotate(
        employee_count=Count('employee')
    ).order_by('name')

    report_from_str = request.GET.get('report_from')
    report_to_str = request.GET.get('report_to')
    report_department_id = request.GET.get('report_department', '')

    report_from = date(today.year, today.month, 1)
    report_to = date(today.year, today.month, calendar.monthrange(today.year, today.month)[1])

    if report_from_str:
        try:
            report_from = datetime.strptime(report_from_str, '%Y-%m-%d').date()
        except ValueError:
            pass

    if report_to_str:
        try:
            report_to = datetime.strptime(report_to_str, '%Y-%m-%d').date()
        except ValueError:
            pass

    report_scope = Employee.objects.select_related('user', 'department')
    if report_department_id:
        report_scope = report_scope.filter(department_id=report_department_id)

    report_scope_ids = list(report_scope.values_list('id', flat=True))

    report_records = Attendance.objects.filter(
        date__range=(report_from, report_to)
    ).select_related(
        'employee__user', 'employee__department', 'added_by'
    )

    if report_scope_ids:
        report_records = report_records.filter(employee_id__in=report_scope_ids)

    report_present_count = report_records.filter(status='Present').count()
    report_late_count = report_records.filter(status='Late').count()
    report_leave_count = report_records.filter(status='Leave').count()
    report_absent_count = max(len(report_scope_ids) - report_records.values('employee_id').distinct().count(), 0)

    report_overtime_count = 0
    for record in report_records:
        if record.check_in and record.check_out:
            check_in_dt = datetime.combine(record.date, record.check_in)
            check_out_dt = datetime.combine(record.date, record.check_out)
            if check_out_dt > check_in_dt + timedelta(hours=8):
                report_overtime_count += 1

    report_total_count = report_records.count()
    # -------------------------------
    # DYNAMIC REPORT DATA
    # -------------------------------

    monthly_summary_rows = []

    for emp in report_scope:
        emp_records = report_records.filter(
            employee=emp
        )

        present = emp_records.filter(
            status='Present'
        ).count()

        late = emp_records.filter(
            status='Late'
        ).count()

        leave = emp_records.filter(
            status='Leave'
        ).count()

        absent = 0

        total_days = (report_to - report_from).days + 1

        attended_days = emp_records.values(
            'date'
        ).distinct().count()

        if total_days > attended_days:
            absent = total_days - attended_days

        total_hours = 0
        total_ot = 0

        for rec in emp_records:
            if rec.check_in and rec.check_out:
                start = datetime.combine(
                    rec.date,
                    rec.check_in
                )
                end = datetime.combine(
                    rec.date,
                    rec.check_out
                )
                diff = end - start
                hours = diff.total_seconds() / 3600
                total_hours += hours

                if hours > 8:
                    total_ot += hours - 8

        monthly_summary_rows.append({
            "employee": emp,
            "present": present,
            "late": late,
            "leave": leave,
            "absent": absent,
            "working_hours": round(total_hours, 2),
            "overtime_hours": round(total_ot, 2),
        })

    # Absent Report Data

    absent_rows = []

    for row in monthly_summary_rows:
        if row["absent"] > 0:
            absent_rows.append(row)

    # Overtime Report Data

    overtime_rows = []

    for record in report_records:
        if record.check_in and record.check_out:
            start = datetime.combine(
                record.date,
                record.check_in
            )
            end = datetime.combine(
                record.date,
                record.check_out
            )
            diff = end - start
            if diff > timedelta(hours=8):
                overtime_rows.append({
                    "employee": record.employee,
                    "date": record.date,
                    "check_in": record.check_in,
                    "check_out": record.check_out,
                    "overtime": round(
                        (diff.total_seconds()/3600)-8,
                        2
                    )
                })

    # -------------------------------
    # Detailed Attendance Report Rows
    # -------------------------------

    report_rows = []

    for record in report_records:
        worked_hours = ""
        overtime_hours = ""

        if record.check_in and record.check_out:
            check_in_dt = datetime.combine(
                record.date,
                record.check_in
            )
            check_out_dt = datetime.combine(
                record.date,
                record.check_out
            )
            diff = check_out_dt - check_in_dt
            worked_hours = round(
                diff.total_seconds() / 3600,
                2
            )
            if diff > timedelta(hours=8):
                overtime_hours = round(
                    (diff - timedelta(hours=8)).total_seconds() / 3600,
                    2
                )

        if record.source == "web":
            source_label = "Web Check-In"
        elif record.source == "manual":
            source_label = "Manual Entry"
        elif record.source == "sync":
            source_label = "Sync Punch Log"
        else:
            source_label = "-"

        report_rows.append({
            "employee": record.employee,
            "date": record.date,
            "check_in": record.check_in,
            "check_out": record.check_out,
            "status": record.status,
            "worked_hours": worked_hours,
            "overtime_hours": overtime_hours,
            "source": source_label,
        })

    # Calendar
    cal = calendar.Calendar(firstweekday=6)

    month_days = []

    for week in cal.monthdayscalendar(
        year,
        month
    ):
        month_days.append([
            {
                'day': day,
                'is_current': (
                    day == today.day
                    and month == today.month
                    and year == today.year
                ),
                'has_record': (
                    day != 0
                    and Attendance.objects.filter(
                        date=date(
                            year,
                            month,
                            day
                        )
                    ).exists()
                )
            }
            for day in week
        ])
    return render(
        request,
        'hrm/attendance.html',

        _ctx(
            request,

            records=records,

            punch_logs=punch_logs,

            manual_entries=manual_entries,

            device_names=device_names,

            devices_online=devices_online,

            punches_synced=punches_synced,

            missing_punches=missing_punches,

            last_sync_label=last_sync_label,


            today_attendance=today_attendance,

            current_employee=current_employee,


            month=month,

            year=year,


            present_count=present_count,

            late_count=late_count,

            on_leave_count=on_leave_count,

            absent_count=absent_count,


            attendance_rows=attendance_rows,


            months=list(range(1, 13)),


            today=today,


            view_mode=view_mode,


            month_days=month_days,


            month_name=calendar.month_name[month],


            active_tab=active_tab,


            employees=employees,


            report_departments=report_departments,

            report_from=report_from,

            report_to=report_to,

            report_department_id=report_department_id,

            report_records=report_records,

            report_total_count=report_total_count,

            report_present_count=report_present_count,

            report_late_count=report_late_count,

            report_leave_count=report_leave_count,

            report_absent_count=report_absent_count,

            report_overtime_count=report_overtime_count,


            monthly_summary_rows=monthly_summary_rows,


            absent_rows=absent_rows,


            overtime_rows=overtime_rows,

            report_rows=report_rows,
            report_type=report_type,

            show_absent_report=show_absent_report,

            show_overtime_report=show_overtime_report,

            show_summary_report=show_summary_report,
        )
    )
@login_required
def attendance_export(request):

    if not has_permission(request.user, "attendance", "view"):
        messages.error(request, "You don't have permission to access Attendance.")
        return redirect("attendance_list")
    today = date.today()
    export_format = request.GET.get('format')

    report_from_str = request.GET.get('report_from')
    report_to_str = request.GET.get('report_to')
    report_department_id = request.GET.get('report_department')
    report_type = request.GET.get('report_type', 'summary')

    if report_from_str or report_to_str:
        if report_from_str:
            try:
                report_from = datetime.strptime(report_from_str, '%Y-%m-%d').date()
            except ValueError:
                report_from = date(today.year, today.month, 1)
        else:
            report_from = date(today.year, today.month, 1)

        if report_to_str:
            try:
                report_to = datetime.strptime(report_to_str, '%Y-%m-%d').date()
            except ValueError:
                report_to = date(today.year, today.month, calendar.monthrange(today.year, today.month)[1])
        else:
            report_to = date(today.year, today.month, calendar.monthrange(today.year, today.month)[1])

        report_scope = Employee.objects.select_related('user', 'department', 'designation')
        if report_department_id:
            report_scope = report_scope.filter(department_id=report_department_id)

        report_scope_ids = list(report_scope.values_list('id', flat=True))

        if report_type == 'absent':
            present_emp_ids = Attendance.objects.filter(
                date__range=(report_from, report_to),
                employee_id__in=report_scope_ids
            ).values_list('employee_id', flat=True).distinct()

            absent_employees = report_scope.exclude(id__in=present_emp_ids)
            rows = [
                [
                    emp.user.get_full_name() or emp.user.username,
                    emp.emp_id,
                    emp.department.name if emp.department else '',
                    emp.designation.title if emp.designation else '',
                    emp.status,
                    emp.join_date.strftime('%Y-%m-%d') if emp.join_date else '',
                ]
                for emp in absent_employees
            ]
            sections = [('Absent Report', ['Employee', 'Employee ID', 'Department', 'Designation', 'Status', 'Join Date'], rows)]
            filename = f'Absent_Report_{report_from}_{report_to}'

        elif report_type == 'overtime':
            records = Attendance.objects.filter(
                date__range=(report_from, report_to),
                employee_id__in=report_scope_ids
            ).select_related('employee__user', 'employee__department').order_by('-date')

            rows = []
            for record in records:
                if record.check_in and record.check_out:
                    check_in_dt = datetime.combine(record.date, record.check_in)
                    check_out_dt = datetime.combine(record.date, record.check_out)
                    diff = check_out_dt - check_in_dt
                    if diff > timedelta(hours=8):
                        hours_worked = round(diff.total_seconds() / 3600.0, 2)
                        overtime_hours = round((diff - timedelta(hours=8)).total_seconds() / 3600.0, 2)
                        rows.append([
                            record.employee.user.get_full_name() or record.employee.user.username,
                            record.employee.emp_id,
                            record.date.strftime('%Y-%m-%d'),
                            record.check_in.strftime('%H:%M'),
                            record.check_out.strftime('%H:%M'),
                            hours_worked,
                            overtime_hours,
                        ])

            sections = [('Overtime Report', ['Employee', 'Employee ID', 'Date', 'Check In', 'Check Out', 'Hours Worked', 'Overtime Hours'], rows)]
            filename = f'Overtime_Report_{report_from}_{report_to}'

        else:
            records = Attendance.objects.filter(
                date__range=(report_from, report_to),
                employee_id__in=report_scope_ids
            ).select_related('employee__user', 'added_by').order_by('-date')

            rows = []
            for record in records:
                rows.append([
                    record.employee.user.get_full_name() or record.employee.user.username,
                    record.employee.emp_id,
                    record.date.strftime('%Y-%m-%d'),
                    record.check_in.strftime('%H:%M') if record.check_in else '',
                    record.check_out.strftime('%H:%M') if record.check_out else '',
                    record.status,
                    record.source,
                    record.added_by.get_full_name() if record.added_by else '',
                    record.note or '',
                ])

            sections = [('Attendance Summary', ['Employee', 'Employee ID', 'Date', 'Check In', 'Check Out', 'Status', 'Source', 'Added By', 'Note'], rows)]
            filename = f'Attendance_Report_{report_from}_{report_to}'

    else:
        month = int(request.GET.get('month', today.month))
        year = int(request.GET.get('year', today.year))

        records = Attendance.objects.filter(
            date__month=month,
            date__year=year,
        ).select_related('employee__user', 'added_by').order_by('-date', 'employee__user__first_name', 'employee__user__last_name')

        rows = []
        for record in records:
            rows.append([
                record.employee.user.get_full_name() or record.employee.user.username,
                record.employee.emp_id,
                record.date.strftime('%Y-%m-%d'),
                record.check_in.strftime('%H:%M') if record.check_in else '',
                record.check_out.strftime('%H:%M') if record.check_out else '',
                record.status,
                record.source,
                record.added_by.get_full_name() if record.added_by else '',
                record.note or '',
            ])

        sections = [('Attendance', ['Employee', 'Employee ID', 'Date', 'Check In', 'Check Out', 'Status', 'Source', 'Added By', 'Note'], rows)]
        filename = f'Attendance_{year}_{month:02d}'

    return _export_tabular_response('Attendance Export', sections, filename, export_format)

# ───────────────────────────────────────────────
# ATTENDANCE CHECK IN
# ───────────────────────────────────────────────

@login_required
def attendance_checkin(request):
    employee = _emp(request)

    if not has_permission(request.user, "attendance", "create"):
        messages.error(request, "You don't have permission to check in.")
        return redirect("attendance_list")
    if not employee:
        return redirect('dashboard')

    today = timezone.now().date()
    current_time = timezone.now().time()
    return_tab = request.POST.get('return_tab', 'today')

    attendance, created = Attendance.objects.get_or_create(
        employee=employee,
        date=today,
        defaults={
            "check_in": current_time,
            "status": "Present",
            "source": "web",
            "added_by": request.user,
        }
    )

    if not created:
        attendance.check_in = current_time
        attendance.status = "Present"
        attendance.source = "web"
        attendance.added_by = request.user
        attendance.save()

    AttendanceLog.objects.create(
        employee=employee,
        date=today,
        time=current_time,
        action="IN",
        device="Web"
    )

    messages.success(
    request,
    f"✅ Checked in at {attendance.check_in.strftime('%I:%M %p')}"
)


    return redirect(f"{reverse('attendance_list')}?tab={return_tab}")

# ───────────────────────────────────────────────
# MANUAL ATTENDANCE ENTRY
# ───────────────────────────────────────────────

@login_required
@manager_or_admin
def attendance_manual(request):

    if not has_permission(request.user, "attendance", "create"):
        messages.error(request, "You don't have permission to add attendance.")
        return redirect("attendance_list")

    if request.method == 'POST':
        employee_id = request.POST.get('employee')
        date_str = request.POST.get('date')
        status = request.POST.get('status')
        check_in = request.POST.get('check_in')
        check_out = request.POST.get('check_out')
        note = request.POST.get('note')
        return_tab = request.POST.get('return_tab', 'manual')
        
        try:
            employee = Employee.objects.get(pk=employee_id)
            attendance_date = datetime.strptime(date_str, '%Y-%m-%d').date()
            
            # Check if attendance already exists for this employee and date
            existing = Attendance.objects.filter(employee=employee, date=attendance_date).first()
            
            if existing:
                # Update existing record
                existing.status = status
                existing.note = note
                existing.source = 'manual'
                existing.added_by = request.user
                if check_in:
                    existing.check_in = datetime.strptime(check_in, '%H:%M').time()
                if check_out:
                    existing.check_out = datetime.strptime(check_out, '%H:%M').time()
                existing.save()
                messages.success(
    request,
    f'Attendance updated for {employee.user.get_full_name()}'
)
            else:
                # Create new record
                attendance = Attendance.objects.create(
                    employee=employee,
                    date=attendance_date,
                    status=status,
                    note=note,
                    source='manual',
                    added_by=request.user,
                )
                if check_in:
                    attendance.check_in = datetime.strptime(check_in, '%H:%M').time()
                if check_out:
                    attendance.check_out = datetime.strptime(check_out, '%H:%M').time()
                attendance.save()
                messages.success(
    request,
    f'Manual attendance added for {employee.user.get_full_name()}'
)
                
        except Employee.DoesNotExist:
            messages.error(request, 'Employee not found')
        except Exception as e:
            messages.error(request, f'Error saving attendance: {str(e)}')
    
    return redirect(f"{reverse('attendance_list')}?tab={return_tab}")






# ───────────────────────────────────────────────
# ATTENDANCE CHECK OUT
# ───────────────────────────────────────────────

@login_required
def attendance_checkout(request):
    if not has_permission(request.user, "attendance", "create"):
        messages.error(request, "You don't have permission to checkout.")
        return redirect("attendance_list")
    employee = _emp(request)


    today = timezone.now().date()
    current_time = timezone.now().time()



    attendance = Attendance.objects.filter(
        employee=employee,
        date=today
    ).first()



    if attendance:

        attendance.check_out = current_time
        attendance.source = 'web'
        attendance.added_by = request.user
        attendance.save()



    AttendanceLog.objects.create(

        employee=employee,
        date=today,
        time=current_time,
        action="OUT",
        device="Web"

    )


    messages.success(
        request,
        "Check Out successful"
    )


    return redirect('attendance_list')


@login_required
def user_guide_faq(request):
    return render(request, 'hrm/guide_faq.html', _ctx(request))


@login_required
def contact_support(request):
    return render(request, 'hrm/contact_support.html', _ctx(request))


@login_required
def dashboard_export(request):
    today = date.today()
    total_employees = Employee.objects.count()
    present = Attendance.objects.filter(date=today, status__in=['Present', 'Late']).count()
    pending_lv = LeaveRequest.objects.filter(from_date__lte=today, to_date__gte=today, status='Approved').count()
    late_today = Attendance.objects.filter(date=today, status='Late').count()
    departments = Department.objects.annotate(employee_count=Count('employee'))

    summary_rows = [
        ['Total Employees', total_employees],
        ['Present Today', present],
        ['On Leave Today', pending_lv],
        ['Late Today', late_today],
    ]
    department_rows = [
        [dept.name, dept.employee_count, round((dept.employee_count / total_employees * 100) if total_employees > 0 else 0, 1)]
        for dept in departments
    ]

    sections = [
        ('Dashboard Summary', ['Metric', 'Value'], summary_rows),
        ('Department Directory', ['Department', 'Employee Count', 'Percentage'], department_rows),
    ]

    return _export_tabular_response('Dashboard Export', sections, 'Dashboard_Summary', request.GET.get('format'))


@login_required
def onboarding_export(request):
    records = OnboardingRecord.objects.select_related('employee__user').all()

    rows = [
        [
            r.employee.user.get_full_name() or r.employee.user.username,
            r.employee.emp_id,
            r.status,
            'Yes' if r.nda_signed else 'No',
            r.nda_signed_at.strftime('%Y-%m-%d %H:%M') if r.nda_signed_at else '',
            r.notes or '',
            r.created_at.strftime('%Y-%m-%d %H:%M') if r.created_at else '',
        ]
        for r in records
    ]

    sections = [
        ('Onboarding Records', ['Employee', 'Employee ID', 'Status', 'NDA Signed?', 'NDA Signed At', 'Notes', 'Created At'], rows),
    ]

    return _export_tabular_response('Onboarding Export', sections, 'Onboarding_Report', request.GET.get('format'))


@login_required
def late_export(request):
    if not has_permission(request.user, "late", "view"):
        messages.error(request, "You don't have permission to access Late Management.")
        return redirect("dashboard")
    if request.user.role == Role.EMPLOYEE:
        late_records = Attendance.objects.filter(
            employee=_emp(request),
            status='Late'
        ).select_related('employee__user', 'late_entry')
    else:
        late_records = Attendance.objects.filter(
            status='Late'
        ).select_related('employee__user', 'late_entry')

    rows = []
    for r in late_records:
        status_label = 'Pending'
        reviewed_by_name = ''
        if r.late_entry:
            if r.late_entry.reviewed_by:
                reviewed_by_name = r.late_entry.reviewed_by.get_full_name() or r.late_entry.reviewed_by.username
                if r.late_entry.excused:
                    status_label = 'Approved'
                else:
                    status_label = 'Rejected'
            elif r.status == 'Approved':
                status_label = 'Approved'

        rows.append([
            r.employee.user.get_full_name() or r.employee.user.username,
            r.employee.emp_id,
            r.date.strftime('%Y-%m-%d'),
            r.check_in.strftime('%H:%M') if r.check_in else '',
            r.late_minutes,
            r.note or '',
            status_label,
            reviewed_by_name,
        ])

    sections = [
        ('Late Entries', ['Employee', 'Employee ID', 'Date', 'Arrival Time', 'Late Minutes', 'Reason', 'Status', 'Reviewed By'], rows),
    ]

    return _export_tabular_response('Late Export', sections, 'Late_Report', request.GET.get('format'))


@login_required
def leave_export(request):
    if not has_permission(request.user, "leave", "view"):
        messages.error(request, "You don't have permission to access Leave Management.")
        return redirect("dashboard")
    emp = _emp(request)
    if request.user.role == Role.EMPLOYEE:
        leaves = LeaveRequest.objects.filter(employee=emp).select_related('employee__user', 'leave_type')
    else:
        leaves = LeaveRequest.objects.select_related('employee__user', 'leave_type').all()

    rows = [
        [
            lv.employee.user.get_full_name() or lv.employee.user.username,
            lv.employee.emp_id,
            lv.leave_type.name,
            lv.from_date.strftime('%Y-%m-%d'),
            lv.to_date.strftime('%Y-%m-%d'),
            (lv.to_date - lv.from_date).days + 1,
            lv.reason or '',
            lv.status,
        ]
        for lv in leaves
    ]

    sections = [
        ('Leave Requests', ['Employee', 'Employee ID', 'Leave Type', 'From Date', 'To Date', 'Days', 'Reason', 'Status'], rows),
    ]

    return _export_tabular_response('Leave Export', sections, 'Leave_Report', request.GET.get('format'))


@manager_or_admin
def punch_log(request):

    logs = AttendanceLog.objects.select_related('employee__user').all()[:100]

    return render(request, 'hrm/punch_log.html', _ctx(request, logs=logs))





# ─── LATE MANAGEMENT ──────────────────────────────────────────────────────────

@login_required
def late_list(request):
    if not has_permission(request.user, "late", "view"):
        messages.error(request, "You don't have permission to access Late Management.")
        return redirect("dashboard")

    from django.utils import timezone


    if request.user.role == Role.EMPLOYEE:

        late = Attendance.objects.filter(
            employee=_emp(request),
            status='Late'
        ).select_related(
            'employee__user',
            'late_entry'
        )

    else:

        late = Attendance.objects.filter(
            status='Late'
        ).select_related(
            'employee__user',
            'late_entry'
        )


    now = timezone.now()


    late_this_month = late.filter(
        date__year=now.year,
        date__month=now.month
    ).count()


    pending_review = late.filter(
        late_entry__reviewed_by__isnull=True
    ).count()


    approved = late.filter(
        late_entry__excused=True
    ).count()


    rejected = late.filter(
        late_entry__excused=False,
        late_entry__reviewed_by__isnull=False
    ).count()



    return render(
        request,
        'hrm/late.html',
        _ctx(
            request,
            late_records=late,
            late_this_month=late_this_month,
            pending_review=pending_review,
            approved=approved,
            rejected=rejected
        )
    )



@login_required
def apply_late(request):

    if not has_permission(request.user, "late", "create"):
        messages.error(request, "You don't have permission to apply late request.")
        return redirect("late_list")


    if request.method == "POST":

        title = request.POST.get("title", "").strip()
        date_raw = request.POST.get("date")
        arrival_raw = request.POST.get("arrival")
        reason = request.POST.get("reason")

        late_date = datetime.strptime(date_raw, "%Y-%m-%d").date()
        arrival_time = datetime.strptime(arrival_raw, "%H:%M").time()

        employee = _emp(request)

        late_minutes = 0
        if employee and employee.shift and employee.shift.start_time:
            scheduled_start = datetime.combine(late_date, employee.shift.start_time)
            arrival_dt = datetime.combine(late_date, arrival_time)
            if arrival_dt > scheduled_start:
                late_minutes = int((arrival_dt - scheduled_start).total_seconds() // 60)


        attendance, _ = Attendance.objects.update_or_create(
            employee=employee,
            date=late_date,
            defaults={
                "status": "Late",
                "check_in": arrival_time,
                "late_minutes": late_minutes,
                "note": reason or "",
                "source": "manual",
                "added_by": request.user,
            }
        )


        LateEntry.objects.update_or_create(
            attendance=attendance,
            defaults={
                "title": title,
                "late_minutes": 0,
                "reason": reason or "",
            }
        )


        messages.success(
            request,
            "Late request submitted successfully."
        )


    return redirect('late_list')



@manager_or_admin
def late_action(request, pk, action):

    if not has_permission(request.user, "late", "edit"):
        messages.error(request, "You don't have permission to update late.")
        return redirect("late_list")


    att = get_object_or_404(
        Attendance,
        pk=pk,
        status='Late'
    )

    le, _ = LateEntry.objects.get_or_create(
        attendance=att,
        defaults={
            'late_minutes': att.late_minutes,
            'reason': att.note or '',
        }
    )

    if action == 'approve':

        le.excused = True
        le.reviewed_by = request.user
        le.reviewed_by_id = request.user.id
        le.save()

        messages.success(
            request,
            f'Late approved for {att.employee}.'
        )

    elif action == 'reject':

        le.excused = False
        le.reviewed_by = request.user
        le.reviewed_by_id = request.user.id
        le.save()

        messages.warning(
            request,
            f'Late rejected for {att.employee}.'
        )

    return redirect('late_list')




@manager_or_admin
def late_review(request, pk):

    if not has_permission(request.user, "late", "edit"):
        messages.error(request, "You don't have permission to review late.")
        return redirect("late_list")



    att = get_object_or_404(
        Attendance,
        pk=pk,
        status='Late'
    )


    le, _ = LateEntry.objects.get_or_create(
        attendance=att,
        defaults={
            'late_minutes': att.late_minutes
        }
    )


    if request.method == 'POST':

        le.warning = request.POST.get(
            'warning',
            'none'
        )

        le.excused = bool(
            request.POST.get('excused')
        )

        le.reason = request.POST.get(
            'reason',
            ''
        )

        le.reviewed_by = request.user

        le.save()


        messages.success(
            request,
            'Late entry reviewed.'
        )


        return redirect(
            'late_list'
        )


    return render(
        request,
        'hrm/late_review.html',
        _ctx(
            request,
            att=att,
            le=le
        )
    )

# ─── LEAVE ────────────────────────────────────────────────────────────────────

@login_required
def leave_list(request):

    if not has_permission(request.user, "leave", "view"):
        messages.error(request, "You don't have permission to access Leave Management.")
        return redirect("dashboard")

    emp = _emp(request)

    if request.user.role == Role.EMPLOYEE:

        leaves = LeaveRequest.objects.filter(
            employee=emp
        )

    else:

        leaves = LeaveRequest.objects.select_related(
            'employee__user',
            'leave_type'
        ).all()


    leave_types = LeaveType.objects.all()


    employees = Employee.objects.select_related(
        'user'
    ).all()


    # ─── Leave Statistics ─────────────────────────────

    today = timezone.now().date()


    on_leave_today = LeaveRequest.objects.filter(
        from_date__lte=today,
        to_date__gte=today,
        status='Approved'
    ).count()


    pending_count = LeaveRequest.objects.filter(
        status='Pending'
    ).count()


    approved_month = LeaveRequest.objects.filter(
        status='Approved',
        from_date__month=today.month,
        from_date__year=today.year
    ).count()


    upcoming = LeaveRequest.objects.filter(
        from_date__gt=today,
        status='Approved'
    ).count()



    return render(
        request,
        'hrm/leave.html',
        _ctx(
            request,

            leaves=leaves,

            leave_types=leave_types,

            employees=employees,


            # stats
            on_leave_today=on_leave_today,
            pending_count=pending_count,
            approved_month=approved_month,
            upcoming=upcoming
        )
    )





@login_required
def leave_apply(request):

    if not has_permission(request.user, "leave", "create"):
        messages.error(request, "You don't have permission to apply for leave.")
        return redirect("leave_list")

    if request.method == "POST":

        title = request.POST.get('title', '').strip()

        employee_id = request.POST.get('employee')

        leave_type_id = request.POST.get('leave_type')


        employee = get_object_or_404(
            Employee,
            id=employee_id
        )


        leave_type = get_object_or_404(
            LeaveType,
            id=leave_type_id
        )


        from_date_raw = request.POST.get('from_date')

        to_date_raw = request.POST.get('to_date')

        reason = request.POST.get('reason')

        attachment = request.FILES.get('attachment')


        half_day = request.POST.get('half_day') == 'yes'



        from_date = datetime.strptime(from_date_raw, '%Y-%m-%d').date()
        to_date = datetime.strptime(to_date_raw, '%Y-%m-%d').date()



        lr = LeaveRequest(

            employee=employee,

            title=title,

            leave_type=leave_type,

            from_date=from_date,

            to_date=to_date,

            reason=reason,

            attachment=attachment,

            half_day=half_day,

            status='Pending'
        )



        if half_day:

            lr.days = 1

        else:

            delta = (to_date - from_date).days + 1

            lr.days = delta



        lr.save()



        messages.success(
            request,
            'Leave request submitted.'
        )


        return redirect('leave_list')



    return redirect('leave_list')





@manager_or_admin
def leave_action(request, pk, action):

    if not has_permission(request.user, "leave", "edit"):
        messages.error(request, "You don't have permission to update leave.")
        return redirect("leave_list")

    lr = get_object_or_404(
        LeaveRequest,
        pk=pk
    )


    if action == 'approve':

        lr.status = 'Approved'

        messages.success(
            request,
            f'Leave approved for {lr.employee}.'
        )


    elif action == 'reject':

        lr.status = 'Rejected'

        messages.warning(
            request,
            f'Leave rejected for {lr.employee}.'
        )


    lr.reviewed_by = request.user

    lr.reviewed_on = timezone.now()

    lr.save()


    return redirect('leave_list')
# ─── TASKS ────────────────────────────────────────────────────────────────────

@login_required
def task_list(request):
    if not has_permission(request.user, "tasks", "view"):
        messages.error(request, "You don't have permission to access Tasks.")
        return redirect("dashboard")

    emp = _emp(request)

    if request.user.role == Role.EMPLOYEE:
        tasks = Task.objects.filter(
            assignee=emp
        ).select_related(
            'assignee__user',
            'project'
        )
    else:
        tasks = Task.objects.select_related(
            'assignee__user',
            'project'
        ).all()

    projects = Project.objects.all()

    if request.user.role == Role.EMPLOYEE and emp:
        employees = Employee.objects.filter(pk=emp.pk).select_related('user')
    else:
        employees = Employee.objects.select_related('user').all()

    statuses = [
        ('To Do', '#F3F4F6', '#374151'),
        ('In Progress', '#FEF3C7', '#D97706'),
        ('Done', '#D1FAE5', '#10B981'),
    ]

    view_mode = request.GET.get('view', 'card')

    context = _ctx(
        request,
        tasks=tasks,
        projects=projects,
        employees=employees,
        statuses=statuses,
        view_mode=view_mode,
        todo_count=tasks.filter(status='To Do').count(),
        progress_count=tasks.filter(status='In Progress').count(),
        done_count=tasks.filter(status='Completed').count(),
    )

    return render(request, 'hrm/tasks.html', context)


@login_required
def task_export(request):
    if not has_permission(request.user, "tasks", "view"):
        messages.error(request, "You don't have permission to export tasks.")
        return redirect("task_list")

    if request.user.role == Role.EMPLOYEE:
        tasks = Task.objects.filter(assignee=_emp(request)).select_related('assignee__user', 'project')
    else:
        tasks = Task.objects.select_related('assignee__user', 'project').all()

    rows = [
        [
            task.title,
            task.project.name,
            task.assignee.user.get_full_name() or task.assignee.user.username,
            task.priority,
            task.status,
            task.progress,
            task.due_date.strftime('%Y-%m-%d') if task.due_date else '',
            task.description or '',
        ]
        for task in tasks
    ]

    sections = [
        ('Tasks', ['Title', 'Project', 'Assignee', 'Priority', 'Status', 'Progress', 'Due Date', 'Description'], rows),
    ]

    return _export_tabular_response('Task Export', sections, 'Tasks', request.GET.get('format'))


@login_required
def task_create(request):
    if not has_permission(request.user, "tasks", "create"):
        messages.error(request, "You don't have permission to create tasks.")
        return redirect("task_list")

    form = TaskForm(request.POST or None)

    if request.method == "GET":
        if request.user.role == Role.EMPLOYEE:
            form.fields['assignee'].initial = _emp(request)

    if form.is_valid():
        task = form.save(commit=False)

        requested_status = request.POST.get('status', '').strip()
        allowed_statuses = {choice for choice, _ in Task.STATUS_CHOICES}
        if requested_status in allowed_statuses:
            task.status = requested_status
        if requested_status == 'In Progress' and (task.progress or 0) == 0:
            task.progress = 20
        elif requested_status == 'Completed':
            task.progress = 100
        elif requested_status == 'To Do':
            task.progress = 0

        task.save()
        messages.success(request, 'Task created.')
        return redirect('task_list')

    return render(
        request,
        'hrm/form.html',
        _ctx(request, form=form, title='New Task', back='task_list')
    )


@login_required
def task_update_status(request, pk):
    if not has_permission(request.user, "tasks", "edit"):
        return JsonResponse({'ok': False, 'error': 'Permission denied.'}, status=403)

    if request.method != 'POST':
        return HttpResponseForbidden('Invalid request method.')

    task = get_object_or_404(Task, pk=pk)

    if request.user.role == Role.EMPLOYEE and task.assignee != _emp(request):
        raise PermissionDenied

    status = request.POST.get('status', '').strip()
    allowed_statuses = {choice for choice, _ in Task.STATUS_CHOICES}

    if status not in allowed_statuses:
        return JsonResponse({'ok': False, 'error': 'Invalid status.'}, status=400)

    task.status = status
    if status == 'In Progress' and task.progress < 50:
        task.progress = 50
    elif status == 'Completed':
        task.progress = 100
    elif status == 'To Do' and task.progress > 0:
        task.progress = 0
    task.save(update_fields=['status', 'progress'])

    return JsonResponse({
        'ok': True,
        'task_id': task.pk,
        'status': task.status,
        'progress': task.progress,
    })


@login_required
def task_edit(request, pk):
    if not has_permission(request.user, "tasks", "edit"):
        messages.error(request, "You don't have permission to edit tasks.")
        return redirect("task_list")

    task = get_object_or_404(Task, pk=pk)

    if request.user.role == Role.EMPLOYEE and task.assignee != _emp(request):
        raise PermissionDenied

    form = TaskForm(request.POST or None, instance=task)

    if form.is_valid():
        form.save()
        messages.success(request, 'Task updated.')
        return redirect('task_list')

    return render(
        request,
        'hrm/form.html',
        _ctx(request, form=form, title='Edit Task', back='task_list')
    )


@login_required
def task_delete(request, pk):
    if not has_permission(request.user, "tasks", "delete"):
        messages.error(request, "You don't have permission to delete tasks.")
        return redirect("task_list")

    task = get_object_or_404(Task, pk=pk)

    user_role = getattr(request.user, 'role', None)
    is_admin = user_role in ['admin', 'super_admin', 'Admin', 'Super Admin']
    if hasattr(Role, 'SUPER_ADMIN'):
        is_admin = is_admin or user_role == Role.SUPER_ADMIN

    if not is_admin:
        emp = _emp(request)
        if not emp or task.assignee != emp:
            raise PermissionDenied("You can only delete your own tasks.")

    task.delete()
    messages.success(request, 'Task deleted.')
    return redirect('task_list')


@login_required
@require_POST
def task_update_color(request, pk):
    task = get_object_or_404(Task, pk=pk)

    if request.user.role == Role.EMPLOYEE and task.assignee != _emp(request):
        return JsonResponse({'success': False, 'error': 'Permission denied'}, status=403)

    color = request.POST.get('color') or request.POST.get('task_color')
    if not color:
        return JsonResponse({'success': False, 'error': 'No color provided'}, status=400)

    task.color = color
    task.save(update_fields=['color'])
    return JsonResponse({'success': True, 'color': task.color})


# ─── NEW: detail (redirects to list + opens modal) ────────────────────────────
@login_required
def task_detail(request, pk):
    if not has_permission(request.user, "tasks", "view"):
        messages.error(request, "You don't have permission to view tasks.")
        return redirect("task_list")

    task = get_object_or_404(Task, pk=pk)

    if request.user.role == Role.EMPLOYEE and task.assignee != _emp(request):
        raise PermissionDenied

    return redirect(f"{reverse('task_list')}?open={task.pk}")


# ─── NEW: 5-step progress (each step = +20%) ─────────────────────────────────
@login_required
@require_POST
def task_update_progress(request, pk):
    if not has_permission(request.user, "tasks", "edit"):
        return JsonResponse({'ok': False, 'error': 'Permission denied.'}, status=403)

    task = get_object_or_404(Task, pk=pk)

    if request.user.role == Role.EMPLOYEE and task.assignee != _emp(request):
        return JsonResponse({'ok': False, 'error': 'Permission denied.'}, status=403)

    try:
        step = int(request.POST.get('step', 0))
    except (TypeError, ValueError):
        return JsonResponse({'ok': False, 'error': 'Invalid step.'}, status=400)

    if step < 0 or step > 5:
        return JsonResponse({'ok': False, 'error': 'Step must be 0-5.'}, status=400)

    # Each step = exactly 20%
    new_progress = step * 20
    task.progress = new_progress

    if new_progress == 0:
        task.status = 'To Do'
    elif new_progress < 100:
        task.status = 'In Progress'
    else:
        task.status = 'Completed'

    task.save(update_fields=['progress', 'status'])

    return JsonResponse({
        'ok': True,
        'progress': task.progress,
        'status': task.status,
    })


# ─── NEW: column color (applies to whole column) ──────────────────────────────
@login_required
@require_POST
def task_update_column_color(request):
    if not has_permission(request.user, "tasks", "edit"):
        return JsonResponse({'ok': False, 'error': 'Permission denied.'}, status=403)

    status = request.POST.get('status', '').strip()
    color  = request.POST.get('color', '').strip()

    allowed = {'To Do', 'In Progress', 'Completed'}
    if status not in allowed or not color:
        return JsonResponse({'ok': False, 'error': 'Invalid data.'}, status=400)

    field_map = {
        'To Do': 'todo_color',
        'In Progress': 'in_progress_color',
        'Completed': 'completed_color',
    }
    field = field_map[status]

    Project.objects.all().update(**{field: color})

    return JsonResponse({'ok': True, 'status': status, 'color': color})
# ─── DOCUMENTS ────────────────────────────────────────────────────────────────

@login_required
def document_list(request):
    if not has_permission(request.user, "documents", "view"):
        messages.error(request, "You don't have permission to access Documents.")
        return redirect("dashboard")

    emp = _emp(request)
    document_type_filter = request.GET.get('document_type', 'all')
    
    print(f"DEBUG: document_type_filter = {document_type_filter}")
    
    # Base query based on user role
    if request.user.role == Role.EMPLOYEE:
        # Employee sees public documents and their own documents
        docs = Document.objects.filter(
            Q(is_public=True) | Q(employee=emp)
        )
    elif request.user.role == Role.MANAGER:
        # Manager sees only their own documents and public documents
        docs = Document.objects.filter(
            Q(is_public=True) | Q(employee=emp)
        )
    else:
        # Super admin sees all documents
        docs = Document.objects.all()
    
    print(f"DEBUG: docs count before filter = {docs.count()}")
    
    # Filter by document type if specified
    if document_type_filter != 'all' and document_type_filter != 'request':
        docs = docs.filter(document_type=document_type_filter)
        print(f"DEBUG: docs count after filter = {docs.count()}")
    
    # Get document requests separately if filtering by request type
    requests = []
    if document_type_filter == 'request':
        requests = DocumentRequest.objects.select_related(
            'requester', 'employee__user'
        ).all().order_by('-requested_at')



    category_count = docs.values(
        'category'
    ).distinct().count()



    total_documents = docs.count()



    pending_requests = docs.filter(
        expiry_date__isnull=False
    ).count()



    generated_documents = docs.filter(
        is_public=True
    ).count()



    return render(
        request,
        'hrm/documents.html',
        _ctx(
            request,
            docs=docs,
            requests=requests,
            document_type_filter=document_type_filter,
            category_count=category_count,
            total_documents=total_documents,
            pending_requests=pending_requests,
            generated_documents=generated_documents
        )
    )


@login_required
def document_export(request):
    if not has_permission(request.user, "documents", "export"):
        messages.error(request, "You don't have permission to export documents.")
        return redirect("document_list")
     

    if request.user.role == Role.EMPLOYEE:
        docs = Document.objects.filter(
            Q(is_public=True) | Q(employee=_emp(request))
        )
    else:
        docs = Document.objects.all()

    rows = [
        [
            doc.name,
            doc.category,
            doc.owner.get_full_name() or doc.owner.username,
            doc.uploaded_at.strftime('%Y-%m-%d %H:%M') if doc.uploaded_at else '',
            doc.expiry_date.strftime('%Y-%m-%d') if doc.expiry_date else '',
            'Yes' if doc.is_public else 'No',
        ]
        for doc in docs
    ]

    sections = [
        ('Documents', ['Name', 'Category', 'Owner', 'Uploaded At', 'Expiry Date', 'Public'], rows),
    ]

    return _export_tabular_response('Document Export', sections, 'Documents', request.GET.get('format'))



@login_required

def document_upload(request):
    form = DocumentForm(request.POST or None, request.FILES or None)

    if form.is_valid():

        d = form.save(commit=False)

        d.owner = request.user
        
        # Set default document_type to 'office' if not specified
        if not d.document_type:
            d.document_type = 'office'

        d.save()

        messages.success(request, 'Document uploaded.')

        return redirect('document_list')

    return render(request, 'hrm/form.html', _ctx(request, form=form, title='Upload Document', back='document_list'))





@admin_required

def document_delete(request, pk):
    if not has_permission(request.user, "documents", "delete"):
        messages.error(request, "You don't have permission to delete documents.")
        return redirect("document_list")

    get_object_or_404(Document, pk=pk).delete()

    messages.success(request, 'Document deleted.')

    return redirect('document_list')


from django.http import FileResponse, Http404
from django.shortcuts import get_object_or_404

@login_required
def document_download(request, pk):
    if not has_permission(request.user, "documents", "view"):
        messages.error(request, "You don't have permission to download documents.")
        return redirect("document_list")

    doc = get_object_or_404(Document, pk=pk)

    # Same visibility rules as document_list
    if request.user.role == Role.EMPLOYEE:
        if not (doc.is_public or doc.employee == _emp(request)):
            raise Http404

    if not doc.file:
        raise Http404("File not found")

    # Use the original filename stored in the FileField
    filename = doc.file.name.split('/')[-1]   # strips the 'documents/' prefix

    return FileResponse(
        doc.file.open('rb'),
        as_attachment=True,
        filename=filename
    )

# ─── DOCUMENTS ────────────────────────────────────────────────────────────────

@login_required
def document_list(request):
    if not has_permission(request.user, "documents", "view"):
        messages.error(request, "You don't have permission to access Documents.")
        return redirect("dashboard")

    emp = _emp(request)
    document_type_filter = request.GET.get('document_type', 'all')

    docs = Document.objects.select_related('owner', 'employee__user').all()

    if request.user.role == Role.EMPLOYEE and emp:
        docs = docs.filter(Q(employee=emp) | Q(is_public=True) | Q(owner=request.user))

    if document_type_filter == 'generated':
        generated_ids = DocumentRequest.objects.filter(
            status='Completed',
            generated_document__isnull=False
        ).values_list('generated_document_id', flat=True)
        docs = docs.filter(pk__in=generated_ids)
    elif document_type_filter == 'employee':
        docs = docs.filter(category__in=['Personal', 'Employee', 'HR'])
    elif document_type_filter == 'office':
        docs = docs.filter(category__in=['Office', 'Policy', 'Company', 'General'])

    if request.user.role == Role.EMPLOYEE:
        requests_qs = DocumentRequest.objects.filter(requester=request.user)
    else:
        requests_qs = DocumentRequest.objects.all()

    requests_qs = requests_qs.select_related(
        'requester', 'employee__user', 'reviewed_by', 'generated_document'
    ).order_by('-requested_at')

    pending_requests = DocumentRequest.objects.filter(status='Pending').count()
    generated_documents = DocumentRequest.objects.filter(status='Completed').count()
    category_count = docs.values('category').distinct().count()
    total_documents = docs.count()

    return render(
        request,
        'hrm/documents.html',
        _ctx(
            request,
            docs=docs,
            requests=requests_qs,
            document_type_filter=document_type_filter,
            category_count=category_count,
            total_documents=total_documents,
            pending_requests=pending_requests,
            generated_documents=generated_documents,
        )
    )


@login_required
def document_export(request):
    if not has_permission(request.user, "documents", "view"):
        messages.error(request, "You don't have permission to export documents.")
        return redirect("document_list")

    docs = Document.objects.select_related('owner', 'employee__user').all()
    rows = [
        [
            d.name,
            d.category or '',
            getattr(d, 'size_display', '') or '',
            d.owner.get_full_name() if d.owner else '',
            d.uploaded_at.strftime('%Y-%m-%d') if d.uploaded_at else '',
        ]
        for d in docs
    ]
    sections = [('Documents', ['Name', 'Category', 'Size', 'Owner', 'Date'], rows)]
    return _export_tabular_response('Document Export', sections, 'Documents', request.GET.get('format'))


@login_required
def document_upload(request):
    if request.user.role not in [Role.SUPER_ADMIN, Role.MANAGER, 'super_admin', 'manager', 'admin']:
        if not has_permission(request.user, "documents", "create"):
            messages.error(request, "You don't have permission to upload documents.")
            return redirect("document_list")

    if request.method == 'POST':
        form = DocumentForm(request.POST, request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.owner = request.user
            doc.save()
            messages.success(request, 'Document uploaded.')
            return redirect('document_list')
    else:
        form = DocumentForm()

    return render(
        request,
        'hrm/form.html',
        _ctx(request, form=form, title='Upload Document', back='document_list')
    )


@login_required
def document_download(request, pk):
    if not has_permission(request.user, "documents", "view"):
        messages.error(request, "You don't have permission to download documents.")
        return redirect("document_list")

    doc = get_object_or_404(Document, pk=pk)
    if not doc.file:
        messages.error(request, 'File not found.')
        return redirect('document_list')

    response = HttpResponse(doc.file.read(), content_type='application/octet-stream')
    response['Content-Disposition'] = f'attachment; filename="{doc.file.name.split("/")[-1]}"'
    return response


@login_required
def document_delete(request, pk):
    if request.user.role not in [Role.SUPER_ADMIN, 'super_admin', 'admin']:
        if not has_permission(request.user, "documents", "delete"):
            messages.error(request, "You don't have permission to delete documents.")
            return redirect("document_list")

    doc = get_object_or_404(Document, pk=pk)
    doc.delete()
    messages.success(request, 'Document deleted.')
    return redirect('document_list')


# ─── DOCUMENT REQUESTS ─────────────────────────────────────────────────────────

@login_required
def document_request_list(request):
    if request.user.role not in [Role.MANAGER, Role.EMPLOYEE, Role.SUPER_ADMIN]:
        if not has_permission(request.user, "documents", "view"):
            messages.error(request, "You don't have permission to view document requests.")
            return redirect("document_list")

    if request.user.role == Role.EMPLOYEE:
        requests = DocumentRequest.objects.filter(requester=request.user).select_related(
            'employee__user', 'reviewed_by'
        )
    else:
        requests = DocumentRequest.objects.select_related(
            'requester', 'employee__user', 'reviewed_by'
        ).all()

    pending_count = DocumentRequest.objects.filter(status='Pending').count()
    generated_count = DocumentRequest.objects.filter(status='Completed').count()

    return render(
        request,
        'hrm/document_requests.html',
        _ctx(
            request,
            requests=requests,
            pending_count=pending_count,
            generated_count=generated_count,
            doc_types=DocumentRequest.DOC_TYPE_CHOICES,
        ),
    )


@login_required
def document_request_create(request):
    if request.user.role not in [Role.MANAGER, Role.EMPLOYEE, Role.SUPER_ADMIN]:
        messages.error(request, "You don't have permission to request documents.")
        return redirect("document_list")

    if request.user.role == Role.EMPLOYEE:
        if not has_permission(request.user, "documents", "create"):
            messages.error(request, "You don't have permission to request documents.")
            return redirect("document_list")

    emp = _emp(request)

    if request.method == 'POST':
        document_type = request.POST.get('document_type')
        custom_type = request.POST.get('custom_type', '')
        reason = request.POST.get('reason')
        employee_id = request.POST.get('employee')

        if request.user.role in [Role.MANAGER, Role.EMPLOYEE]:
            employee = emp
        else:
            employee = get_object_or_404(Employee, pk=employee_id)

        DocumentRequest.objects.create(
            requester=request.user,
            employee=employee,
            document_type=document_type,
            custom_type=custom_type,
            reason=reason,
            status='Pending',
        )
        messages.success(request, 'Document request submitted successfully.')
        return redirect('document_request_list')

    if request.user.role in [Role.MANAGER, Role.EMPLOYEE]:
        employees = [emp] if emp else []
    else:
        employees = Employee.objects.select_related('user').all()

    return render(
        request,
        'hrm/document_request_form.html',
        _ctx(request, employees=employees, doc_types=DocumentRequest.DOC_TYPE_CHOICES),
    )


@manager_or_admin
def document_request_approve(request, pk):
    if not has_permission(request.user, "documents", "edit"):
        messages.error(request, "You don't have permission to approve document requests.")
        return redirect("document_request_list")

    doc_request = get_object_or_404(DocumentRequest, pk=pk)

    if request.method == 'POST':
        action = request.POST.get('action')
        notes = request.POST.get('notes', '')

        doc_request.reviewed_by = request.user
        doc_request.reviewed_at = timezone.now()
        doc_request.notes = notes

        if action == 'approve':
            doc_request.status = 'Approved'
            messages.success(request, f'Document request for {doc_request.employee.user.get_full_name()} approved.')
        elif action == 'reject':
            doc_request.status = 'Rejected'
            messages.warning(request, f'Document request for {doc_request.employee.user.get_full_name()} rejected.')

        doc_request.save()
        return redirect('document_request_list')

    return render(
        request,
        'hrm/document_request_approve.html',
        _ctx(request, doc_request=doc_request),
    )


@manager_or_admin
def document_generate(request, pk):
    if not has_permission(request.user, "documents", "create"):
        messages.error(request, "You don't have permission to generate documents.")
        return redirect("document_request_list")

    doc_request = get_object_or_404(DocumentRequest, pk=pk)

    if request.method == 'POST':
        uploaded_file = request.FILES.get('generated_file')
        document_name = request.POST.get('document_name')
        category = request.POST.get('category', 'Personal')

        if uploaded_file:
            document = Document.objects.create(
                name=document_name or f"{doc_request.document_type} - {doc_request.employee.user.get_full_name()}",
                file=uploaded_file,
                category=category,
                owner=request.user,
                employee=doc_request.employee,
                is_public=False,
            )
            doc_request.generated_document = document
            doc_request.status = 'Completed'
            doc_request.reviewed_by = request.user
            doc_request.reviewed_at = timezone.now()
            doc_request.save()
            messages.success(request, 'Document generated and uploaded successfully.')
            return redirect('document_request_list')
        else:
            messages.error(request, 'Please upload the generated document.')

    return render(
        request,
        'hrm/document_generate.html',
        _ctx(request, doc_request=doc_request, categories=Document.CAT_CHOICES),
    )

# ─── SALARY / PAYROLL ─────────────────────────────────────────────────────────
@login_required
def salary_list(request):
    if not has_permission(request.user, "payroll", "view"):
        messages.error(request, "You don't have permission to access Payroll.")
        return redirect("dashboard")


    emp = _emp(request)

    if request.user.role == Role.EMPLOYEE:

        payslips = (
            Payslip.objects
            .filter(employee=emp)
            .select_related(
                'employee__user',
                'structure'
            )
        )

    else:

        payslips = (
            Payslip.objects
            .select_related(
                'employee__user',
                'structure'
            )
            .all()
        )


    employees = Employee.objects.select_related(
        'user'
    ).filter(status='Active')


    gross_payroll = Decimal("0")
    total_bonus = Decimal("0")
    total_deductions = Decimal("0")
    net_disbursed = Decimal("0")


    for ps in payslips:

        gross_payroll += ps.structure.gross

        total_bonus += ps.bonus

        total_deductions += (
            ps.structure.tax_deduction +
            ps.structure.pf_deduction +
            ps.deduction
        )

        net_disbursed += ps.net_pay


    adjustments = PayrollAdjustment.objects.select_related(
        'employee__user',
        'created_by'
    )[:10]


    return render(

        request,

        "hrm/salary_list.html",

        _ctx(

            request,

            payslips=payslips,

            employees=employees,

            adjustments=adjustments,

            gross_payroll=gross_payroll,

            total_bonus=total_bonus,

            total_deductions=total_deductions,

            net_disbursed=net_disbursed,

        )

    )


@login_required
def salary_export(request):
    if not has_permission(request.user, "payroll", "view"):
        messages.error(request, "You don't have permission to export payroll.")
        return redirect("salary_list")

    emp = _emp(request)

    if request.user.role == Role.EMPLOYEE:
        payslips = (
            Payslip.objects
            .filter(employee=emp)
            .select_related('employee__user', 'structure')
        )
        adjustments = PayrollAdjustment.objects.filter(employee=emp).select_related('employee__user', 'created_by')
    else:
        payslips = (
            Payslip.objects
            .select_related('employee__user', 'structure')
            .all()
        )
        adjustments = PayrollAdjustment.objects.select_related('employee__user', 'created_by').all()

    payslip_rows = [
        [
            ps.employee.user.get_full_name() or ps.employee.user.username,
            f'{ps.month:02d}/{ps.year}',
            ps.status,
            ps.structure.gross,
            ps.bonus,
            ps.structure.tax_deduction + ps.structure.pf_deduction + ps.deduction,
            ps.net_pay,
        ]
        for ps in payslips
    ]
    adjustment_rows = [
        [
            adj.employee.user.get_full_name() or adj.employee.user.username,
            adj.adjustment_type,
            adj.amount,
            adj.reason,
            adj.created_by.get_full_name() if adj.created_by else '',
            adj.created_at.strftime('%Y-%m-%d %H:%M') if adj.created_at else '',
        ]
        for adj in adjustments
    ]

    sections = [
        ('Payslips', ['Employee', 'Period', 'Status', 'Gross', 'Bonus', 'Total Deductions', 'Net Pay'], payslip_rows),
        ('Adjustments', ['Employee', 'Type', 'Amount', 'Reason', 'Created By', 'Created At'], adjustment_rows),
    ]

    return _export_tabular_response('Salary Export', sections, 'Salary_Report', request.GET.get('format'))
@login_required
def salary_adjustments(request):
    if not has_permission(request.user, "payroll", "view"):
        messages.error(request, "You don't have permission to view adjustments.")
        return redirect("salary_list")


    adjustments = PayrollAdjustment.objects.select_related(
        'employee__user',
        'created_by'
    )

    return render(
        request,
        'hrm/salary_adjustments.html',
        _ctx(
            request,
            adjustments=adjustments
        )
    )


@login_required
def employee_breakdown(request):
    if not has_permission(request.user, "payroll", "view"):
        messages.error(request, "You don't have permission to view salary breakdown.")
        return redirect("salary_list")

    employees = Employee.objects.select_related(
        'user'
    ).filter(status='Active')

    selected_employee = None
    structure = None
    latest_payslip = None

    employee_id = request.GET.get("employee")

    if employee_id:

        selected_employee = Employee.objects.filter(
            id=employee_id
        ).first()

        if selected_employee:

            structure = SalaryStructure.objects.filter(
                employee=selected_employee,
                is_active=True
            ).first()

            latest_payslip = Payslip.objects.filter(
                employee=selected_employee
            ).order_by(
                "-year",
                "-month"
            ).first()

    return render(
        request,
        "hrm/employee_breakdown.html",
        _ctx(
            request,
            employees=employees,
            selected_employee=selected_employee,
            structure=structure,
            latest_payslip=latest_payslip,
        )
    )


@login_required
def salary_history(request):
    if not has_permission(request.user, "payroll", "view"):
        messages.error(request, "You don't have permission to view salary history.")
        return redirect("salary_list")


    payslips = Payslip.objects.select_related(
        "employee__user",
        "structure"
    )

    return render(
        request,
        "hrm/salary_history.html",
        _ctx(
            request,
            payslips=payslips,
        )
    )


@login_required
def salary_reports(request):
    if not has_permission(request.user, "payroll", "view"):
        messages.error(request, "You don't have permission to view salary reports.")
        return redirect("salary_list")

    return render(
        request,
        "hrm/salary_reports.html",
        _ctx(request)
    )


@admin_required
def salary_structure_create(request):
    if not has_permission(request.user, "payroll", "create"):
        messages.error(request, "You don't have permission to create salary structure.")
        return redirect("salary_list")

    return render(
        request,
        "hrm/salary_structure.html",
        _ctx(request)
    )


@admin_required
def run_payroll(request):
    if not has_permission(request.user, "payroll", "create"):
        messages.error(request, "You don't have permission to run payroll.")
        return redirect("salary_list")

    return render(
        request,
        "hrm/run_payroll.html",
        _ctx(request)
    )
# ─── PETTY CASH ───────────────────────────────────────────────────────────────

@admin_required

def petty_cash(request):
    if not has_permission(request.user, "petty_cash", "view"):
        messages.error(request, "You don't have permission to access Petty Cash.")
        return redirect("dashboard")


    entries = PettyCashLedger.objects.select_related('created_by').all()

    total_credit = entries.filter(entry_type='Credit').aggregate(s=Sum('amount'))['s'] or 0

    total_debit  = entries.filter(entry_type='Debit').aggregate(s=Sum('amount'))['s'] or 0

    balance = total_credit - total_debit

    return render(request, 'hrm/petty_cash.html', _ctx(

        request, entries=entries, total_credit=total_credit,

        total_debit=total_debit, balance=balance,

    ))


@admin_required
def petty_cash_export(request):
    if not has_permission(request.user, "petty_cash", "view"):
        messages.error(request, "You don't have permission to export petty cash.")
        return redirect("petty_cash")

    entries = PettyCashLedger.objects.select_related('created_by').all()

    rows = [
        [
            entry.date.strftime('%Y-%m-%d'),
            entry.description,
            entry.category,
            entry.entry_type,
            entry.amount,
            entry.balance,
            entry.created_by.get_full_name() or entry.created_by.username,
            entry.note,
        ]
        for entry in entries
    ]

    sections = [
        ('Petty Cash Ledger', ['Date', 'Description', 'Category', 'Type', 'Amount', 'Balance', 'Created By', 'Note'], rows),
    ]

    return _export_tabular_response('Petty Cash Export', sections, 'Petty_Cash_Ledger', request.GET.get('format'))





@admin_required

def petty_cash_add(request):
    if not has_permission(request.user, "petty_cash", "create"):
        messages.error(request, "You don't have permission to add petty cash entry.")
        return redirect("petty_cash")


    form = PettyCashForm(request.POST or None)

    if form.is_valid():

        entry = form.save(commit=False)

        entry.created_by = request.user

        # compute running balance

        last = PettyCashLedger.objects.order_by('-id').first()

        prev_bal = last.balance if last else 0

        if entry.entry_type == 'Credit':

            entry.balance = prev_bal + entry.amount

        else:

            entry.balance = prev_bal - entry.amount

        entry.save()

        messages.success(request, 'Entry added.')

        return redirect('petty_cash')

    return render(request, 'hrm/form.html', _ctx(request, form=form, title='Add Cash Entry', back='petty_cash'))





# ─── ASSETS ───────────────────────────────────────────────────────────────────

@login_required

def asset_list(request):
    if not has_permission(request.user, "assets", "view"):
        messages.error(request, "You don't have permission to access Assets.")
        return redirect("dashboard")

    if request.user.role == Role.EMPLOYEE:

        assets = Asset.objects.filter(assigned_to=_emp(request))

    else:

        assets = Asset.objects.select_related('assigned_to__user').all()

    return render(request, 'hrm/assets.html', _ctx(request, assets=assets))


@login_required
def asset_export(request):
    if not has_permission(request.user, "assets", "view"):
        messages.error(request, "You don't have permission to export assets.")
        return redirect("asset_list")

    if request.user.role == Role.EMPLOYEE:
        assets = Asset.objects.filter(assigned_to=_emp(request)).select_related('assigned_to__user')
    else:
        assets = Asset.objects.select_related('assigned_to__user').all()

    rows = [
        [
            asset.asset_id,
            asset.name,
            asset.asset_type,
            asset.serial_no,
            asset.assigned_to.user.get_full_name() if asset.assigned_to else '',
            asset.assigned_on.strftime('%Y-%m-%d') if asset.assigned_on else '',
            asset.status,
            asset.purchase_date.strftime('%Y-%m-%d') if asset.purchase_date else '',
            asset.purchase_cost or '',
            asset.note,
        ]
        for asset in assets
    ]

    sections = [
        ('Assets', ['Asset ID', 'Name', 'Type', 'Serial No', 'Assigned To', 'Assigned On', 'Status', 'Purchase Date', 'Cost', 'Note'], rows),
    ]

    return _export_tabular_response('Asset Export', sections, 'Assets', request.GET.get('format'))





@manager_or_admin

def asset_create(request):
    if not has_permission(request.user, "assets", "create"):
        messages.error(request, "You don't have permission to create assets.")
        return redirect("asset_list")


    form = AssetForm(request.POST or None)

    if form.is_valid():

        form.save()

        messages.success(request, 'Asset added.')

        return redirect('asset_list')

    return render(request, 'hrm/form.html', _ctx(request, form=form, title='Add Asset', back='asset_list'))





@manager_or_admin

def asset_edit(request, pk):

    if not has_permission(request.user, "assets", "edit"):
        messages.error(request, "You don't have permission to edit assets.")
        return redirect("asset_list")

    asset = get_object_or_404(Asset, pk=pk)

    form  = AssetForm(request.POST or None, instance=asset)

    if form.is_valid():

        form.save()

        messages.success(request, 'Asset updated.')

        return redirect('asset_list')

    return render(request, 'hrm/form.html', _ctx(request, form=form, title='Edit Asset', back='asset_list'))





@admin_required

def asset_delete(request, pk):
    if not has_permission(request.user, "assets", "delete"):
        messages.error(request, "You don't have permission to delete assets.")
        return redirect("asset_list")

    get_object_or_404(Asset, pk=pk).delete()

    messages.success(request, 'Asset deleted.')

    return redirect('asset_list')





# ─── FILES & CREDENTIALS ──────────────────────────────────────────────────────

@admin_required

def files_list(request):

    if not has_permission(request.user, "files", "view"):
        messages.error(request, "You don't have permission to access files.")
        return redirect("dashboard")

    files = SecureFile.objects.filter(owner=request.user)

    return render(request, 'hrm/files.html', _ctx(request, files=files))


@admin_required
def files_export(request):
    if not has_permission(request.user, "files", "view"):
        messages.error(request, "You don't have permission to export files.")
        return redirect("files_list")

    files = SecureFile.objects.filter(owner=request.user)

    rows = [
        [
            file.name,
            file.owner.get_full_name() or file.owner.username,
            file.uploaded_at.strftime('%Y-%m-%d %H:%M') if file.uploaded_at else '',
            file.note,
        ]
        for file in files
    ]

    sections = [
        ('Secure Files', ['Name', 'Owner', 'Uploaded At', 'Note'], rows),
    ]

    return _export_tabular_response('Files Export', sections, 'Secure_Files', request.GET.get('format'))





@admin_required

def file_upload(request):
    if not has_permission(request.user, "files", "create"):
        messages.error(request, "You don't have permission to upload files.")
        return redirect("files_list")

    if request.method == 'POST':

        name = request.POST.get('name', '')

        f    = request.FILES.get('file')

        note = request.POST.get('note', '')

        if f:

            SecureFile.objects.create(name=name, file=f, note=note, owner=request.user)

            messages.success(request, 'File uploaded securely.')

        return redirect('files_list')

    return render(request, 'hrm/file_upload.html', _ctx(request))




# ─── ONBOARDING ───────────────────────────────────────────────────────────────

@manager_or_admin
def onboarding_list(request):
    if not has_permission(request.user, "onboarding", "view"):
        messages.error(request, "You don't have permission to access Onboarding.")
        return redirect("dashboard")

    records = OnboardingRecord.objects.select_related(
        'employee__user'
    ).all()

    employees = Employee.objects.select_related(
        'user',
        'department',
        'designation'
    ).all()

    return render(
        request,
        'hrm/onboarding.html',
        _ctx(
            request,
            records=records,
            employees=employees
        )
    )


@manager_or_admin
def onboarding_sign_nda(request, pk):
    if not has_permission(request.user, "onboarding", "edit"):
        messages.error(request, "You don't have permission to update onboarding.")
        return redirect("onboarding_list")

    rec = get_object_or_404(
        OnboardingRecord,
        pk=pk
    )

    rec.nda_signed = True
    rec.nda_signed_at = timezone.now()
    rec.status = 'NDA Signed'
    rec.save()

    messages.success(
        request,
        f'NDA marked as signed for {rec.employee}.'
    )

    return redirect('onboarding_list')


@manager_or_admin
def onboarding_complete(request):
    if not has_permission(request.user, "onboarding", "create"):
        messages.error(request, "You don't have permission to complete onboarding.")
        return redirect("onboarding_list")

    if request.method == "POST":

        employee_id = request.POST.get('employee_id')

        employee = get_object_or_404(
            Employee,
            id=employee_id
        )

        rec, created = OnboardingRecord.objects.get_or_create(
            employee=employee
        )

        rec.nda_signed = True
        rec.nda_signed_at = timezone.now()
        rec.status = "Completed"
        rec.save()

        messages.success(
            request,
            f"{employee} onboarding completed successfully."
        )

    return redirect('onboarding_list')

# ─── CONFIGURATION (admin only) ───────────────────────────────────────────────

@login_required

def config(request):
    if not has_permission(request.user, "configuration", "view"):
        messages.error(request, "You don't have permission to access Configuration.")
        return redirect("dashboard")

    departments = Department.objects.annotate(emp_count=Count('employee'))

    designations = Designation.objects.select_related('department').all()

    shifts       = Shift.objects.all()

    holidays     = Holiday.objects.all()

    email_tpls   = EmailTemplate.objects.all()

    notif_rules  = NotificationRule.objects.all()

    banks        = BankAccount.objects.select_related('employee__user').all()



    forms = {

        'dept_form':  DepartmentForm(),

        'desig_form': DesignationForm(),

        'shift_form': ShiftForm(),

        'hol_form':   HolidayForm(),

        'notif_form': NotificationRuleForm(),

    }

    return render(request, 'hrm/config.html', _ctx(

        request,

        departments=departments, designations=designations, shifts=shifts,

        holidays=holidays, email_tpls=email_tpls, notif_rules=notif_rules,

        banks=banks, **forms,

    ))
@login_required
def dept_create(request):
    if not has_permission(request.user, "configuration", "create"):
        messages.error(request, "You don't have permission to create department.")
        return redirect("config")

    if request.method == "POST":
        form = DepartmentForm(request.POST)

        if form.is_valid():
            form.save()
            messages.success(
                request,
                "Department created successfully"
            )
            return redirect('department_list')

    else:
        form = DepartmentForm()

    return render(
        request,
        'hrm_app/dept_form.html',
        {
            'form': form
        }
    )

@login_required
def site_settings(request):

    if not has_permission(request.user, "configuration", "edit"):
        messages.error(request, "You don't have permission to update settings.")
        return redirect("dashboard")

    settings, created = SiteSettings.objects.get_or_create(id=1)

    if request.method == "POST":

        form = SiteSettingsForm(
            request.POST,
            request.FILES,
            instance=settings
        )

        if form.is_valid():
            form.save()
            messages.success(request, "Site settings saved successfully.")
            return redirect("site_settings")

    else:
        form = SiteSettingsForm(instance=settings)

    return render(
        request,
        "hrm/site_settings.html",
        {
            "form": form,
            "settings": settings,
        },
    )


@login_required

def desig_create(request):
    if not has_permission(request.user, "configuration", "create"):
        messages.error(request, "You don't have permission to create department.")
        return redirect("config")

    form = DesignationForm(request.POST)

    if form.is_valid():

        form.save()

        messages.success(request, 'Designation added.')

    return redirect('config')





@admin_required

def shift_create(request):
    if not has_permission(request.user, "configuration", "create"):
        messages.error(request, "You don't have permission to create shift.")
        return redirect("config")

    form = ShiftForm(request.POST)

    if form.is_valid():

        form.save()

        messages.success(request, 'Shift added.')

    return redirect('config')




@login_required
def holiday_create(request):

    if not has_permission(request.user, "configuration", "create"):
        messages.error(request, "You don't have permission to create holiday.")
        return redirect("config")

    form = HolidayForm(request.POST)

    if form.is_valid():

        form.save()

        messages.success(request, 'Holiday added.')

    return redirect('config')

    form = HolidayForm(request.POST)

    if form.is_valid():

        form.save()

        messages.success(request, 'Holiday added.')

    return redirect('config')





@admin_required

def notif_rule_toggle(request, pk):
    if not has_permission(request.user, "configuration", "edit"):
        messages.error(request, "You don't have permission to edit notification rules.")
        return redirect("config")


    rule = get_object_or_404(NotificationRule, pk=pk)

    rule.is_active = not rule.is_active

    rule.save()

    messages.success(request, f'Notification rule {"enabled" if rule.is_active else "disabled"}.')

    return redirect('config')





# ─── PERMISSIONS PAGE ─────────────────────────────────────────────────────────
@admin_required
def permissions_view(request):

    modules = UserPermission.MODULE_CHOICES

    permissions = []

    for module_key, module_name in modules:
        permissions.append({
            "module": module_key,
            "name": module_name,
        })

    return render(
        request,
        "hrm/permissions.html",
        _ctx(
            request,
            permissions=permissions
        )
    )


@login_required
def calendar_view(request):
    if not has_permission(request.user, "calendar", "view"):
        messages.error(request, "You don't have permission to access Calendar.")
        return redirect("dashboard")
    today = date.today()
    month = int(request.GET.get('month', today.month))
    year  = int(request.GET.get('year', today.year))

    # Get approved leaves for the month
    leaves = LeaveRequest.objects.filter(
        status='Approved',
        from_date__year=year,
        from_date__month=month
    ).select_related('employee__user')

    # Get holidays
    holidays = Holiday.objects.filter(date__year=year, date__month=month)

    cal = calendar.Calendar(firstweekday=6)
    month_days = []
    for week in cal.monthdayscalendar(year, month):
        month_days.append([
            {
                'day': day,
                'is_current': day == today.day and month == today.month and year == today.year,
                'has_leave': any(
                    l.from_date.day <= day <= l.to_date.day
                    for l in leaves
                    if l.from_date.month == month and l.from_date.year == year
                )
            }
            for day in week
        ])

    return render(request, 'hrm/calendar.html', _ctx(
        request,
        month=month,
        year=year,
        month_name=calendar.month_name[month],
        month_days=month_days,
    ))


# ─── SUPPORT PAGES ─────────────────────────────────────────────────────────────



@login_required
def contact_support(request):
    return render(request, 'hrm/contact_support.html', _ctx(request))


